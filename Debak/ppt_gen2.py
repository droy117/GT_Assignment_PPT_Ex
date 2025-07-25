import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from tkinter import filedialog, messagebox
import customtkinter as ctk
import matplotlib.pyplot as plt
import os
import io
import asyncio
import requests
import re
import json
from dotenv import load_dotenv
import threading

load_dotenv()


# --- Helper function to truncate long labels ---
def truncate_label(label, length=20):
    """Truncates a label if it's longer than the specified length."""
    label_str = str(label)  # Ensure the label is a string
    if len(label_str) > length:
        return label_str[:length] + '...'
    return label_str

# --- Constants and Theme Settings ---
APP_TITLE = "PPTEx"
WINDOW_SIZE = "1200x800"
COLUMN_ACTIONS = [
    "Ignore",
    "Group Slides by this Column",
    "Summarize as Bullet Points",
    "Create Bar Chart",
    "Create Pie Chart",
    "Create Histogram",
    "Create Line Chart",
    "Include in Data Table"
]
AI_CALL_BATCH_SIZE = 3 # Process 3 columns per API call
MAX_CATEGORIES_FOR_CHARTS = 9 # For performance, group smaller categories into "Others"

# --- Main Application Class ---
class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        # --- Window Setup ---
        self.title(APP_TITLE)
        self.geometry(WINDOW_SIZE)
        ctk.set_appearance_mode("Dark")
        ctk.set_default_color_theme("blue")

        # --- Class Variables ---
        self.dataframe = None
        self.original_df = None
        self.data_file_path = ""
        self.template_path = ""
        self.column_widgets = {}
        self.max_score_info = {}
        self.is_cancelling = False # For cancellable generation
        self.generation_thread = None # To hold the generation thread

        # --- Set default template path if it exists ---
        default_template_file = './Templates/Template_1_default.pptx'
        if os.path.exists(default_template_file):
            self.template_path = default_template_file
        else:
            print(f"Default template not found at {default_template_file}. Will use a blank presentation.")

        # --- Main Layout Configuration ---
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(2, weight=1) # The scrollable frame gets the most space

        # --- Top Frame for File Selection ---
        self.file_frame = ctk.CTkFrame(self, corner_radius=10)
        self.file_frame.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        self.file_frame.grid_columnconfigure(1, weight=1)

        self.browse_data_button = ctk.CTkButton(self.file_frame, text="Select Data File (Excel, CSV)", command=self.load_data_file)
        self.browse_data_button.grid(row=0, column=0, padx=10, pady=10)
        self.data_file_label = ctk.CTkLabel(self.file_frame, text="No data file selected", anchor="w")
        self.data_file_label.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        self.browse_template_button = ctk.CTkButton(self.file_frame, text="Select PowerPoint Template (Optional)", command=self.load_template_file)
        self.browse_template_button.grid(row=1, column=0, padx=10, pady=10)
        self.template_label = ctk.CTkLabel(self.file_frame, text=f"Template: {os.path.basename(self.template_path)}" if self.template_path else "No template selected. A default will be used.", anchor="w")
        self.template_label.grid(row=1, column=1, padx=10, pady=10, sticky="ew")

        # --- Middle Frame for Settings ---
        self.settings_frame = ctk.CTkFrame(self, corner_radius=10)
        self.settings_frame.grid(row=1, column=0, padx=20, pady=10, sticky="ew")
        self.settings_frame.grid_columnconfigure(1, weight=1)

        ctk.CTkLabel(self.settings_frame, text="Presentation Title:").grid(row=0, column=0, padx=(10,0), pady=10, sticky="w")
        self.ppt_title_entry = ctk.CTkEntry(self.settings_frame, placeholder_text="e.g., Website Assessment Report")
        self.ppt_title_entry.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        # --- Main Scrollable Frame for Column Mappings ---
        self.scrollable_frame = ctk.CTkScrollableFrame(self, label_text="Map Your Data Columns to PowerPoint Elements")
        self.scrollable_frame.grid(row=2, column=0, padx=20, pady=10, sticky="nsew")
        self.scrollable_frame.grid_columnconfigure(1, weight=1)

        # --- Bottom Frame for Actions ---
        self.bottom_frame = ctk.CTkFrame(self, corner_radius=10)
        self.bottom_frame.grid(row=3, column=0, padx=20, pady=(10, 0), sticky="ew")
        self.bottom_frame.grid_columnconfigure(0, weight=1)

        self.generate_button = ctk.CTkButton(self.bottom_frame, text="Generate Presentation", command=self.start_generation_thread, state="disabled")
        self.generate_button.grid(row=0, column=0, padx=10, pady=10, sticky="ew")

        # --- Status Log Textbox ---
        self.status_log = ctk.CTkTextbox(self, height=150, state="disabled", wrap="word")
        self.status_log.grid(row=4, column=0, padx=20, pady=(10, 20), sticky="ew")

        # Configure color tags for the log
        self.status_log.tag_config("INFO", foreground="gray")
        self.status_log.tag_config("SUCCESS", foreground="green")
        self.status_log.tag_config("ERROR", foreground="red")
        self.status_log.tag_config("WARN", foreground="orange")
        self.status_log.tag_config("CYAN", foreground="cyan")
        self.status_log.tag_config("WHITE", foreground="white")

        self.log_status("Ready. Please select a data file.", "INFO")

    def log_status(self, message, level="INFO"):
        """Logs a message to the status text box with appropriate color and auto-scrolls."""
        self.status_log.configure(state="normal")
        self.status_log.insert("end", f"{message}\n", level.upper())
        self.status_log.configure(state="disabled")
        self.status_log.see("end") # Auto-scroll to the bottom
        self.update_idletasks() # Ensure UI updates immediately

    def clean_data(self, df):
        """
        Cleans specific columns to make them chart-friendly.
        REVISED to eliminate pandas UserWarning.
        """
        self.log_status("Cleaning data...", "INFO")
        cleaned_df = df.copy()
        self.max_score_info = {}

        for col in cleaned_df.columns:
            # Skip if the column is already a numeric type
            if pd.api.types.is_numeric_dtype(cleaned_df[col].dtype):
                continue

            # Process non-numeric columns that contain some data
            if cleaned_df[col].notna().any():
                col_as_str = cleaned_df[col].astype(str)

                # --- FIX: Directly extract instead of using .contains() first ---
                pattern = r'^\s*([0-9.]+)\s*(out of|/)\s*([0-9.]+)'
                extracted_data = col_as_str.str.extract(pattern)

                # Check if any rows matched and yielded valid data in the first capture group
                if not extracted_data[0].isnull().all():
                    scores = pd.to_numeric(extracted_data[0], errors='coerce')
                    max_vals = pd.to_numeric(extracted_data[2], errors='coerce')

                    if scores.notna().any():
                        new_col_name = f"{col} (Score)"
                        cleaned_df[new_col_name] = scores
                        first_valid_max = max_vals.dropna().iloc[0] if max_vals.notna().any() else None
                        if first_valid_max:
                            self.max_score_info[new_col_name] = first_valid_max
                        cleaned_df.drop(columns=[col], inplace=True)
                        self.log_status(f"Cleaned and extracted scores from '{col}' into '{new_col_name}'.", "INFO")
                        continue # Move to the next column

                # Attempt to convert object columns with numbers-as-strings (e.g., "1,000")
                if pd.api.types.is_object_dtype(cleaned_df[col].dtype):
                    cleaned_series = pd.to_numeric(
                        col_as_str.str.replace(',', '', regex=False),
                        errors='coerce'
                    )
                    # If conversion is successful for any value, apply it
                    if cleaned_series.notna().any():
                        cleaned_df[col] = cleaned_series

        self.log_status("Data cleaning complete.", "SUCCESS")
        return cleaned_df

    def load_data_file(self):
        self.data_file_path = filedialog.askopenfilename(
            filetypes=(("All Data Files", "*.xlsx *.xls *.csv *.tsv"),("Excel files", "*.xlsx *.xls"),("CSV files", "*.csv"),("TSV files", "*.tsv"),("All files", "*.*"))
        )
        if not self.data_file_path:
            self.log_status("File selection cancelled.", "WARN")
            return

        self.data_file_label.configure(text=os.path.basename(self.data_file_path))
        self.log_status(f"Reading data file: {os.path.basename(self.data_file_path)}", "WHITE")

        try:
            file_ext = os.path.splitext(self.data_file_path)[1].lower()
            header_row = self._find_header_row(self.data_file_path)
            self.log_status(f"Detected header at row {header_row + 1}.", "INFO")

            if file_ext in ['.xlsx', '.xls']:
                df = pd.read_excel(self.data_file_path, header=header_row)
            elif file_ext == '.csv':
                df = pd.read_csv(self.data_file_path, header=header_row)
            elif file_ext == '.tsv':
                df = pd.read_csv(self.data_file_path, sep='\t', header=header_row)
            else:
                raise ValueError("Unsupported file type.")

            df.dropna(how='all', inplace=True)
            df = df.loc[:, ~df.columns.str.contains('^Unnamed')]
            self.original_df = df.copy()
            self.dataframe = self.clean_data(df)
            self.populate_column_mappings()
            self.generate_button.configure(state="normal")
            self.log_status("File loaded successfully. Please map columns.", "SUCCESS")
        except Exception as e:
            messagebox.showerror("Error Loading File", f"Could not read the data file.\n\nError: {e}")
            self.log_status(f"Error loading file: {e}", "ERROR")
            self.dataframe = None
            self.generate_button.configure(state="disabled")

    def _find_header_row(self, file_path):
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext in ['.xlsx', '.xls']:
                temp_df = pd.read_excel(file_path, header=None, nrows=10)
            else:
                temp_df = pd.read_csv(file_path, header=None, nrows=10, sep=None, engine='python')
            for i, row in temp_df.iterrows():
                if row.notna().sum() > len(row) / 2 and all(isinstance(x, str) for x in row if pd.notna(x)):
                    if row.nunique() >= len(row) / 2:
                        return i
        except Exception:
            return 0
        return 0

    def load_template_file(self):
        self.template_path = filedialog.askopenfilename(filetypes=(("PowerPoint templates", "*.pptx"),))
        if self.template_path:
            self.template_label.configure(text=os.path.basename(self.template_path))
            self.log_status(f"Template selected: {os.path.basename(self.template_path)}", "INFO")
        else:
            self.template_label.configure(text="No template selected. A default will be used.")
            self.log_status("Template selection cancelled. Using default.", "WARN")

    def populate_column_mappings(self):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        self.column_widgets.clear()
        if self.dataframe is None: return

        self.log_status("Populating column mappings with intelligent defaults...", "INFO")
        for i, col_name in enumerate(self.dataframe.columns):
            label = ctk.CTkLabel(self.scrollable_frame, text=col_name, wraplength=250)
            label.grid(row=i, column=0, padx=10, pady=(5, 10), sticky="w")
            unique_count = self.dataframe[col_name].nunique()
            dtype = self.dataframe[col_name].dtype
            col_lower = col_name.lower()
            if col_name in self.max_score_info:
                default_action = "Create Bar Chart"
            elif "risk" in col_lower or "level" in col_lower:
                 default_action = "Create Pie Chart"
            elif dtype in ['int64', 'float64']:
                default_action = "Create Histogram"
            elif 0 < unique_count < 6 and dtype == 'object':
                default_action = "Create Pie Chart"
            elif dtype == 'object' and unique_count > len(self.dataframe) * 0.8:
                 default_action = "Summarize as Bullet Points"
            else:
                default_action = "Ignore"
            dropdown_var = ctk.StringVar(value=default_action)
            dropdown = ctk.CTkOptionMenu(self.scrollable_frame, values=COLUMN_ACTIONS, variable=dropdown_var, width=250)
            dropdown.grid(row=i, column=1, padx=10, pady=(5, 10), sticky="e")
            self.column_widgets[col_name] = [label, dropdown]
        self.log_status("Column mapping populated.", "INFO")


    # --- PPTX HELPER FUNCTIONS ---
    def add_title_slide(self, prs, slide_title, slide_subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_subtitle

    def add_chart_and_insight_slide(self, prs, title, chart_image_stream, insight_text):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        placeholder_left = slide.placeholders[1]
        ph_element = placeholder_left.element
        ph_element.getparent().remove(ph_element)
        chart_image_stream.seek(0)
        slide.shapes.add_picture(chart_image_stream, placeholder_left.left, placeholder_left.top, width=placeholder_left.width)
        placeholder_right = slide.placeholders[2]
        tf = placeholder_right.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        lines = insight_text.strip().replace('**', '').split('\n')
        for line in lines:
            line = line.strip()
            if not line: continue
            if line.endswith(':'):
                p = tf.add_paragraph()
                p.text = line
                p.level = 0; p.font.bold = True; p.font.size = Pt(18)
            elif line.startswith(('* ', '- ')):
                p = tf.add_paragraph()
                p.text = line[2:]
                p.level = 1; p.font.size = Pt(14)
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0; p.font.size = Pt(14)

    def add_section_header_slide(self, prs, title):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        placeholders_to_remove = [shp for shp in slide.shapes if shp.is_placeholder and shp.shape_id != title_shape.shape_id]
        for shape in placeholders_to_remove:
            sp_element = shape.element
            sp_element.getparent().remove(sp_element)

    def add_bullet_point_slide(self, prs, title, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = str(point)
            p.level = 0

    def add_table_slide(self, prs, title, df_table):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        rows, cols = df_table.shape[0] + 1, df_table.shape[1]
        if len(slide.placeholders) < 2:
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(5.5))
        else:
            placeholder = slide.placeholders[1]
            table_shape = slide.shapes.add_table(rows, cols, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
            placeholder.element.getparent().remove(placeholder.element)
        table = table_shape.table
        for c, col_name in enumerate(df_table.columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for r_idx, row_data in enumerate(df_table.itertuples(index=False), start=1):
            for c_idx, cell_data in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(cell_data)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                if (r_idx % 2) == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(220, 230, 241)

    def summarize_categorical_data(self, series: pd.Series, threshold: int):
        """Groups smaller categories into 'Others' for cleaner charts and better performance."""
        if not pd.api.types.is_object_dtype(series.dtype) and not pd.api.types.is_categorical_dtype(series.dtype):
            return series

        counts = series.value_counts()
        if len(counts) > threshold:
            top_n = counts.head(threshold).copy() # Use copy to avoid SettingWithCopyWarning
            other_count = counts.iloc[threshold:].sum()

            num_other_cats = len(counts) - threshold
            other_label = f"Others ({num_other_cats} categories)"

            top_n[other_label] = other_count
            return top_n
        return counts

    async def get_ai_insights_for_batch(self, column_batch: list):
        """
        Handles raw data series and pre-summarized value counts using a 'hint'.
        """
        batch_col_names = [col_name for col_name, _, _ in column_batch]
        self.log_status(f"Generating AI insights for: {', '.join(batch_col_names)}...", "CYAN")

        data_strings_concatenated = ""
        for col_name, data, hint in column_batch:
            data_summary = ""
            if hint == 'counts':
                data_summary = data.to_string()
            elif hint == 'series':
                if pd.api.types.is_numeric_dtype(data.dtype):
                    data_summary = data.describe().to_string()
                else:
                    data_summary = data.value_counts().to_string()

            data_strings_concatenated += f"Column: '{col_name}'\n---\n{data_summary}\n---\n\n"

        prompt = f"""
        As a data analyst, analyze the following data summaries for multiple columns from a dataset.
        The data represents charts that your user will see. For data with an 'Other' category, acknowledge that it represents a collection of smaller groups.

        --- DATA SUMMARIES ---
        {data_strings_concatenated}
        --- END DATA ---

        Based on this data, provide a concise analysis for each column for a PowerPoint slide.
        Your response MUST be a valid JSON object. The keys of the object should be the exact column names provided.
        For each column name key, the value should be a string containing the analysis. Structure the analysis string for each column into two sections using markdown-style headers:

        **Summary Insight:**
        - Write a 2-3 sentence interpretation of what this data reveals. What is the main takeaway?

        **Key Metrics:**
        - List 2-4 key, quantifiable metrics derived from the data.
        - For numerical data, include: Average, Max, Min, and count.
        - For categorical/summarized data, include: The most frequent category and its count, and the least frequent category. If an 'Other' group exists, mention its size.

        Example JSON output format:
        {{
          "Column Name A": "**Summary Insight:**\\n- Brief analysis of A.\\n\\n**Key Metrics:**\\n- Metric 1: Value\\n- Metric 2: Value",
          "Column Name B": "**Summary Insight:**\\n- Brief analysis of B.\\n\\n**Key Metrics:**\\n- Metric 1: Value\\n- Metric 2: Value"
        }}
        """
        try:
            apiKey = os.getenv("GEMINI_API_KEY")
            if not apiKey: return {name: "Error: GEMINI_API_KEY environment variable not found." for name in batch_col_names}

            payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}]}
            apiUrl = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={apiKey}"

            response = await asyncio.to_thread(requests.post, apiUrl, json=payload, timeout=120)
            response.raise_for_status()

            result = response.json()
            if result.get('candidates'):
                json_text = result['candidates'][0]['content']['parts'][0]['text']
                json_text = json_text.strip().lstrip('```json').rstrip('```')

                insights = json.loads(json_text)
                self.log_status(f"Successfully received AI insights for: {', '.join(insights.keys())}", "INFO")
                return insights
            else:
                error_detail = f"Could not generate AI summary. Details: {result.get('promptFeedback', 'No content.')}"
                self.log_status(error_detail, "ERROR")
                return {name: f"Error: {error_detail}" for name in batch_col_names}

        except json.JSONDecodeError as e:
            error_msg = f"AI response was not valid JSON: {e}"
            self.log_status(error_msg, "ERROR")
            return {name: error_msg for name in batch_col_names}
        except Exception as e:
            error_msg = f"An unexpected error occurred during AI insight generation: {e}"
            self.log_status(error_msg, "ERROR")
            return {name: error_msg for name in batch_col_names}


    async def generate_plots_for_df(self, prs, df_subset, group_title=""):
        """
        REVISED: Sets a thread-safe Matplotlib backend to prevent warnings.
        Summarizes large categorical data, fetches AI insights concurrently,
        and checks for cancellation requests.
        """
        # --- FIX: Set a non-GUI backend for Matplotlib to make it thread-safe ---
        import matplotlib
        matplotlib.use('Agg')

        mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}
        chart_actions = ["Create Bar Chart", "Create Pie Chart", "Create Histogram", "Create Line Chart"]

        # --- 1. Pre-compute summarized data for categorical charts ---
        precomputed_counts = {}
        for col, action in mappings.items():
            if self.is_cancelling: return
            if action in ["Create Bar Chart", "Create Pie Chart"]:
                if col in df_subset.columns and pd.api.types.is_object_dtype(df_subset[col].dtype):
                    if df_subset[col].nunique() > 0 and not df_subset[col].isnull().all():
                        summarized = self.summarize_categorical_data(df_subset[col], MAX_CATEGORIES_FOR_CHARTS)
                        precomputed_counts[col] = summarized

        # --- 2. Gather data for AI analysis ---
        columns_for_ai = []
        for col, action in mappings.items():
            if self.is_cancelling: return
            if action in chart_actions and col in df_subset.columns and df_subset[col].nunique() > 0 and not df_subset[col].isnull().all():
                if col in precomputed_counts:
                    columns_for_ai.append((col, precomputed_counts[col], 'counts'))
                else:
                    columns_for_ai.append((col, df_subset[col], 'series'))

        all_insights = {}
        if columns_for_ai:
            if self.is_cancelling: return
            batches = [columns_for_ai[i:i + AI_CALL_BATCH_SIZE] for i in range(0, len(columns_for_ai), AI_CALL_BATCH_SIZE)]
            tasks = [self.get_ai_insights_for_batch(batch) for batch in batches]
            results_list = await asyncio.gather(*tasks)
            for result_dict in results_list:
                if isinstance(result_dict, dict):
                    all_insights.update(result_dict)

        # --- 3. Generate slides ---
        for col, action in mappings.items():
            if self.is_cancelling: return
            if action not in chart_actions: continue
            if col not in df_subset.columns or df_subset[col].nunique() == 0 or df_subset[col].isnull().all(): continue

            self.log_status(f"Creating chart for '{col}'...", "WHITE")

            plt.style.use('seaborn-v0_8-talk')
            fig, ax = plt.subplots(figsize=(10, 6))
            chart_title = f"{group_title}: {col}" if group_title else f"Analysis of {col}"

            # --- Generate Chart Image ---
            try:
                if action == "Create Bar Chart":
                    if col in precomputed_counts:
                        data_counts = precomputed_counts[col]
                        data_counts.index = data_counts.index.map(truncate_label)
                        data_counts.plot(kind='bar', ax=ax, color=plt.cm.viridis.colors)
                        ax.set_ylabel("Count")
                    elif col in self.max_score_info:
                        max_y_val = self.max_score_info[col]
                        df_subset[col].plot(kind='bar', ax=ax, color=plt.cm.viridis(df_subset[col].values / max_y_val))
                        ax.set_ylabel("Score"); ax.set_xlabel("Data Point Index")
                        ax.set_ylim(0, max_y_val * 1.05)
                        chart_title = f"{group_title}: {col.replace(' (Score)', '')}"
                    else:
                         df_subset[col].value_counts().plot(kind='bar', ax=ax, color=plt.cm.viridis.colors)
                         ax.set_ylabel("Count")
                    plt.xticks(rotation=45, ha='right')

                elif action == "Create Pie Chart":
                    data_counts = precomputed_counts.get(col, df_subset[col].value_counts())
                    truncated_labels = data_counts.index.map(truncate_label)
                    ax.pie(data_counts, labels=truncated_labels, autopct='%1.1f%%', startangle=140, colors=plt.cm.Pastel1.colors)
                    ax.axis('equal')

                elif action == "Create Histogram":
                    if pd.api.types.is_numeric_dtype(df_subset[col]):
                        df_subset[col].dropna().plot(kind='hist', ax=ax, bins=15, color='skyblue', ec='black')
                        ax.set_ylabel("Frequency"); ax.set_xlabel(col)
                    else:
                        self.log_status(f"Skipping Histogram for non-numeric column: '{col}'", "WARN")
                        plt.close(fig); continue

                elif action == "Create Line Chart":
                     if pd.api.types.is_numeric_dtype(df_subset[col]):
                        ax.plot(df_subset.index, df_subset[col], marker='o', linestyle='-')
                        ax.set_ylabel(col); ax.set_xlabel("Index")
                     else:
                        self.log_status(f"Skipping Line Chart for non-numeric column: '{col}'", "WARN")
                        plt.close(fig); continue

                plt.tight_layout()
                img_stream = io.BytesIO()
                plt.savefig(img_stream, format='png', dpi=200, bbox_inches='tight')
                plt.close(fig)
                img_stream.seek(0)

                insight_text = all_insights.get(col, "AI insight could not be generated for this chart.")
                self.add_chart_and_insight_slide(prs, chart_title, img_stream, insight_text)
                self.log_status(f"Successfully added slide for '{col}'.", "INFO")

            except Exception as e:
                plt.close(fig)
                self.log_status(f"Failed to create chart for '{col}': {e}", "ERROR")
                continue

        # --- Handle Table Generation ---
        table_cols = [col for col, action in mappings.items() if action == "Include in Data Table"]
        if table_cols:
            if self.is_cancelling: return
            self.log_status("Generating data table slide...", "WHITE")
            table_df = df_subset[table_cols]
            table_title = f"{group_title}: Data Summary" if group_title else "Detailed Data Summary"
            self.add_table_slide(prs, table_title, table_df)
            self.log_status("Data table slide created.", "INFO")

    async def generate_presentation_async(self):
        if self.dataframe is None:
            messagebox.showerror("Error", "No data has been loaded.")
            return

        self.log_status("Starting presentation generation...", "WHITE")

        try:
            prs = Presentation(self.template_path) if self.template_path and os.path.exists(self.template_path) else Presentation()
            if self.is_cancelling: return

            if self.template_path and os.path.exists(self.template_path):
                self.log_status(f"Using template: {os.path.basename(self.template_path)}", "INFO")
                for i in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            else:
                self.log_status("No template found or specified. Using a blank 16:9 presentation.", "INFO")
                prs.slide_width = Inches(16); prs.slide_height = Inches(9)

            if self.is_cancelling: return

            ppt_title = self.ppt_title_entry.get() or "Data Analysis Report"
            subtitle_text = f"Source: {os.path.basename(self.data_file_path)}"
            self.add_title_slide(prs, ppt_title, subtitle_text)
            self.log_status("Added title slide.", "INFO")

            if self.is_cancelling: return

            mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}

            summary_cols = [col for col, action in mappings.items() if action == "Summarize as Bullet Points"]
            if summary_cols:
                self.log_status("Generating summary slide...", "WHITE")
                all_bullets = []
                for col in summary_cols:
                    bullets = self.dataframe[col].dropna().unique().tolist()
                    if bullets: all_bullets.extend(bullets)
                if all_bullets:
                    self.add_bullet_point_slide(prs, "Key Findings & Observations", all_bullets)
                    self.log_status("Summary slide created.", "INFO")

            if self.is_cancelling: return

            grouping_col = next((col for col, action in mappings.items() if action == "Group Slides by this Column"), None)
            if grouping_col:
                self.log_status(f"Grouping slides by column: '{grouping_col}'...", "WHITE")
                unique_groups = self.dataframe[grouping_col].dropna().unique()
                for group in unique_groups:
                    if self.is_cancelling: return
                    self.log_status(f"Generating slides for group: '{group}'...", "WHITE")
                    self.add_section_header_slide(prs, f"Detailed Analysis for: {group}")
                    df_subset = self.dataframe[self.dataframe[grouping_col] == group]
                    await self.generate_plots_for_df(prs, df_subset, group_title=str(group))
            else:
                if self.is_cancelling: return
                self.log_status("Generating slides for the entire dataset...", "WHITE")
                await self.generate_plots_for_df(prs, self.dataframe)

            if self.is_cancelling:
                self.log_status("Generation cancelled before saving.", "WARN")
                return

            self.add_section_header_slide(prs, "Thank You")

            save_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint Presentation", "*.pptx")],
                initialfile=f"{ppt_title.replace(' ', '_')}_Report.pptx",
                title="Save Presentation As"
            )

            if save_path:
                prs.save(save_path)
                self.log_status(f"Success! Presentation saved to {os.path.basename(save_path)}", "SUCCESS")
                messagebox.showinfo("Success", f"Presentation generated successfully!\n\nSaved at: {save_path}")
            else:
                self.log_status("Save operation cancelled by user.", "WARN")

        except Exception as e:
            error_message = f"An unexpected error occurred during generation: {e}"
            self.log_status(error_message, "ERROR")
            messagebox.showerror("Generation Error", error_message, detail=str(e))
            import traceback
            traceback.print_exc()

    # --- Threading and Cancellation Logic ---

    def start_generation_thread(self):
        """Starts the presentation generation in a separate thread to keep the UI responsive."""
        if self.generation_thread and self.generation_thread.is_alive():
            self.log_status("Generation is already in progress.", "WARN")
            return

        self.is_cancelling = False
        # Replace the 'Generate' button with a red 'Cancel' button
        self.generate_button.configure(
            text="Cancel Generation",
            command=self.cancel_generation,
            fg_color="#D32F2F",
            hover_color="#B71C1C"
        )

        self.generation_thread = threading.Thread(target=self._run_async_generation, daemon=True)
        self.generation_thread.start()

    def cancel_generation(self):
        """Signals the generation thread to cancel its operation."""
        if self.generation_thread and self.generation_thread.is_alive():
            self.log_status("Cancellation requested. Will stop after the current step.", "WARN")
            self.is_cancelling = True
            # Disable the button to prevent multiple clicks while waiting for the thread to stop
            self.generate_button.configure(text="Cancelling...", state="disabled")

    def _reset_generate_button_state(self):
        """Resets the generate button to its initial state. Must be called from the main GUI thread."""
        self.generate_button.configure(
            text="Generate Presentation",
            command=self.start_generation_thread,
            state="normal" if self.dataframe is not None else "disabled",
            fg_color=ctk.ThemeManager.theme["CTkButton"]["fg_color"],
            hover_color=ctk.ThemeManager.theme["CTkButton"]["hover_color"]
        )
        self.is_cancelling = False
        self.generation_thread = None

    def _run_async_generation(self):
        """
        The target function for the worker thread. It runs the asyncio event loop for the async generation task.
        """
        try:
            asyncio.run(self.generate_presentation_async())
        except Exception as e:
            # Log any exceptions that occur during the async operation
            self.log_status(f"A runtime error occurred in the generation thread: {e}", "ERROR")
            # Schedule a messagebox to be shown from the main thread for thread safety
            self.after(0, lambda: messagebox.showerror("Runtime Error", f"An unexpected error occurred during generation:\n\n{e}"))
            import traceback
            traceback.print_exc()
        finally:
            # No matter what happens (success, error, or cancellation),
            # schedule the button reset to run on the main GUI thread.
            self.after(0, self._reset_generate_button_state)

if __name__ == "__main__":
    app = App()
    app.mainloop()
