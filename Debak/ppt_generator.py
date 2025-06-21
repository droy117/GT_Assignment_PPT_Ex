import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from tkinter import filedialog, messagebox
import customtkinter as ctk
import matplotlib.pyplot as plt
import os
import io

# --- Constants and Theme Settings ---
APP_TITLE = "Dynamic Excel to PowerPoint Automation"
WINDOW_SIZE = "1000x750"
CHART_TYPES = ["Ignore", "Bar Chart", "Pie Chart", "Include in Table", "Use as Slide Title"]

# --- Main Application Class ---
class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        # --- Window Setup ---
        self.title(APP_TITLE)
        self.geometry(WINDOW_SIZE)
        ctk.set_appearance_mode("Dark")
        ctk.set_default_color_theme("blue")

        # --- Class Variables ---
        self.dataframe = None
        self.excel_path = ""
        self.column_widgets = {} # Stores {column_name: [label_widget, dropdown_widget]}

        # --- Main Layout Configuration ---
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=1)

        # --- Top Frame for File Selection ---
        self.top_frame = ctk.CTkFrame(self, corner_radius=10)
        self.top_frame.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        self.top_frame.grid_columnconfigure(1, weight=1)

        self.browse_button = ctk.CTkButton(self.top_frame, text="Select Excel File", command=self.load_excel_file)
        self.browse_button.grid(row=0, column=0, padx=10, pady=10)

        self.file_label = ctk.CTkLabel(self.top_frame, text="No file selected", anchor="w")
        self.file_label.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        # --- Main Scrollable Frame for Dynamic Column Mappings ---
        self.scrollable_frame = ctk.CTkScrollableFrame(self, label_text="Map Your Data Columns to PowerPoint Elements")
        self.scrollable_frame.grid(row=1, column=0, padx=20, pady=10, sticky="nsew")
        self.scrollable_frame.grid_columnconfigure(1, weight=1)

        # --- Bottom Frame for Actions and Settings ---
        self.bottom_frame = ctk.CTkFrame(self, corner_radius=10)
        self.bottom_frame.grid(row=2, column=0, padx=20, pady=(10, 20), sticky="ew")
        self.bottom_frame.grid_columnconfigure(1, weight=1)

        ctk.CTkLabel(self.bottom_frame, text="Presentation Title:").grid(row=0, column=0, padx=(10,0), pady=10)
        self.ppt_title_entry = ctk.CTkEntry(self.bottom_frame, placeholder_text="e.g., Q3 Website Assessment Report")
        self.ppt_title_entry.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        self.generate_button = ctk.CTkButton(self.bottom_frame, text="Generate Presentation", command=self.generate_presentation, state="disabled")
        self.generate_button.grid(row=1, column=0, columnspan=2, padx=10, pady=10, sticky="ew")

        self.status_label = ctk.CTkLabel(self, text="Status: Ready. Please select an Excel file.", text_color="gray")
        self.status_label.grid(row=3, column=0, padx=20, pady=(0, 10), sticky="w")

    def load_excel_file(self):
        """Opens a dialog to select an Excel file, intelligently finds the header row,
        reads the data, and populates the UI with column mapping options."""
        self.excel_path = filedialog.askopenfilename(
            title="Select an Excel File",
            filetypes=(("Excel files", "*.xlsx *.xls"), ("All files", "*.*"))
        )
        if not self.excel_path:
            return

        self.file_label.configure(text=os.path.basename(self.excel_path))
        self.status_label.configure(text="Status: Reading Excel file...", text_color="white")
        self.update_idletasks() # Force UI update to show status

        try:
            # --- Intelligent Header Detection ---
            xls = pd.ExcelFile(self.excel_path)
            sheet_name = xls.sheet_names[0] # Use the first sheet by default

            # Read the first 10 rows to find the most likely header row
            temp_df = pd.read_excel(xls, sheet_name=sheet_name, header=None, nrows=10)

            header_row_index = 0
            # A good heuristic for a header row is that it has few nulls and its non-null values are all strings.
            for i, row in temp_df.iterrows():
                # Check if more than half the cells are not NA and all of those non-NA cells are strings
                if row.notna().sum() > len(row) / 2 and all(isinstance(x, str) for x in row if pd.notna(x)):
                    header_row_index = i
                    break

            # Read the actual data using the detected header row
            self.dataframe = pd.read_excel(xls, sheet_name=sheet_name, header=header_row_index)
            self.dataframe.dropna(how='all', inplace=True) # Drop fully empty rows

            self.populate_column_mappings()
            self.generate_button.configure(state="normal")
            self.status_label.configure(text="Status: File loaded successfully. Please map columns.", text_color="green")

        except Exception as e:
            messagebox.showerror("Error Loading File", f"Could not read the Excel file.\n\nError: {e}")
            self.status_label.configure(text="Status: Error loading file.", text_color="red")
            self.file_label.configure(text="No file selected")
            self.dataframe = None
            self.generate_button.configure(state="disabled")

    def populate_column_mappings(self):
        """Clears and re-creates UI elements for mapping each Excel column to a report action."""
        # Clear any existing widgets from the scrollable frame
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        self.column_widgets.clear()

        if self.dataframe is None:
            return

        # Create new widgets for each column in the loaded DataFrame
        for i, col_name in enumerate(self.dataframe.columns):
            label = ctk.CTkLabel(self.scrollable_frame, text=col_name)
            label.grid(row=i, column=0, padx=10, pady=(5, 10), sticky="w")

            dropdown = ctk.CTkOptionMenu(self.scrollable_frame, values=CHART_TYPES, width=200)
            dropdown.grid(row=i, column=1, padx=10, pady=(5, 10), sticky="e")

            self.column_widgets[col_name] = [label, dropdown]

    def add_chart_slide(self, prs, title, chart_image_stream):
        """Adds a new slide with a title and a chart image (from a memory stream)."""
        slide_layout = prs.slide_layouts[5] # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        # Add chart image, centered on the slide
        img_width = Inches(8)
        img_height = Inches(5.5)
        left = (prs.slide_width - img_width) / 2
        top = Inches(1.75)
        slide.shapes.add_picture(chart_image_stream, left, top, width=img_width, height=img_height)

    def add_table_slide(self, prs, title, df_table):
        """Adds a new slide with a title and a formatted table from a DataFrame."""
        slide_layout = prs.slide_layouts[5] # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        rows, cols = df_table.shape[0] + 1, df_table.shape[1]
        left, top, width = Inches(1), Inches(2.0), Inches(14)
        height = Inches(0.5) * rows

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Write table headers with styling
        for col_index, col_name in enumerate(df_table.columns):
            cell = table.cell(0, col_index)
            cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 82, 129) # Dark Blue Header

        # Write data rows
        for r_idx, row_data in enumerate(df_table.itertuples(index=False), start=1):
            for c_idx, cell_data in enumerate(row_data):
                table.cell(r_idx, c_idx).text = str(cell_data)

    def generate_presentation(self):
        """Main logic to generate the PowerPoint file based on the user's column mappings."""
        if self.dataframe is None:
            messagebox.showerror("Error", "No Excel data has been loaded.")
            return

        user_ppt_title = self.ppt_title_entry.get() or "Data Analysis Report"

        try:
            self.status_label.configure(text="Status: Generating presentation...", text_color="white")
            self.update_idletasks()

            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)

            # --- Add Title Slide ---
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = user_ppt_title
            subtitle.text = f"Analysis of {os.path.basename(self.excel_path)}"

            # --- Process Mappings and Generate Slides ---
            # Generate Chart Slides
            for col, widgets in self.column_widgets.items():
                choice = widgets[1].get()
                if choice in ["Bar Chart", "Pie Chart"]:
                    plt.style.use('seaborn-v0_8-darkgrid')
                    fig, ax = plt.subplots(figsize=(10, 6))

                    data_counts = self.dataframe[col].value_counts()

                    if choice == "Bar Chart":
                        data_counts.plot(kind='bar', ax=ax)
                        ax.set_ylabel("Count")
                        ax.tick_params(axis='x', rotation=45)
                    elif choice == "Pie Chart":
                        data_counts.plot(kind='pie', ax=ax, autopct='%1.1f%%', startangle=90)
                        ax.set_ylabel('') # Hide y-label for pie charts

                    ax.set_title(f"Distribution of '{col}'", fontsize=16)
                    plt.tight_layout()

                    # Save chart to a memory stream to avoid creating temp files
                    img_stream = io.BytesIO()
                    plt.savefig(img_stream, format='png', dpi=200)
                    plt.close(fig)
                    img_stream.seek(0)

                    self.add_chart_slide(prs, f"Analysis of: {col}", img_stream)

            # Generate Table Slide
            table_cols = [col for col, widgets in self.column_widgets.items() if widgets[1].get() == "Include in Table"]
            if table_cols:
                table_df = self.dataframe[table_cols]
                self.add_table_slide(prs, "Detailed Data Summary", table_df)

            # --- Save the Presentation ---
            save_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint Presentation", "*.pptx")],
                initialfile=f"{user_ppt_title.replace(' ', '_')}.pptx",
                title="Save Presentation As"
            )

            if save_path:
                prs.save(save_path)
                self.status_label.configure(text=f"Status: Success! Saved to {os.path.basename(save_path)}", text_color="green")
                messagebox.showinfo("Success", f"Presentation generated successfully!\n\nSaved at: {save_path}")
            else:
                self.status_label.configure(text="Status: Save operation cancelled.", text_color="orange")

        except Exception as e:
            self.status_label.configure(text="Status: An error occurred during generation.", text_color="red")
            messagebox.showerror("Generation Error", f"An error occurred while creating the presentation:\n\n{e}")

# --- Main Execution ---
if __name__ == "__main__":
    app = App()
    app.mainloop()
