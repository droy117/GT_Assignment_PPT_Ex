import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from tkinter import filedialog, messagebox
import customtkinter as ctk
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
import os
import winreg
import io
import asyncio
import requests

def get_windows_theme():
    """
    Returns 'Light' or 'Dark' depending on the current Windows app theme.
    """
    try:
        registry = winreg.ConnectRegistry(None, winreg.HKEY_CURRENT_USER)
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Themes\Personalize"
        key = winreg.OpenKey(registry, key_path)
        value, _ = winreg.QueryValueEx(key, "AppsUseLightTheme")
        return "Light" if value == 1 else "Dark"
    except Exception:
        return "Light"  # Fallback default

# --- Constants and Theme Settings ---
APP_TITLE = "PPTEx"
WINDOW_SIZE = "1200x800"
COLUMN_ACTIONS = [
    "Ignore",
    "Group Slides by this Column",
    "Use for AI Summary",
    "Summarize as Bullet Points",
    "Create Bar Chart",
    "Create Pie Chart",
    "Include in Data Table"
]

# --- Main Application Class ---
class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        # --- Window Setup ---
        self.title(APP_TITLE)
        self.geometry(WINDOW_SIZE)
        ctk.set_appearance_mode(get_windows_theme())
        ctk.set_default_color_theme("blue")

        # --- Class Variables ---
        self.dataframe = None
        self.data_file_path = ""
        self.template_path = ""
        self.column_widgets = {}

        # --- Main Layout Configuration ---
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(2, weight=1)

        # --- Top Frame for File Selection ---
        self.file_frame = ctk.CTkFrame(self, corner_radius=10)
        self.file_frame.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        self.file_frame.grid_columnconfigure(1, weight=1)

        self.browse_data_button = ctk.CTkButton(self.file_frame, text="Select Data File (Excel, CSV)", command=self.load_data_file)
        self.browse_data_button.grid(row=0, column=0, padx=10, pady=10)
        self.data_file_label = ctk.CTkLabel(self.file_frame, text="No data file selected", anchor="w")
        self.data_file_label.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        self.browse_template_button = ctk.CTkButton(self.file_frame, text="Select PowerPoint Template (Optional)", command=self.load_template_file)
        self.browse_template_button.grid(row=1, column=0, padx=10, pady=10)
        self.template_label = ctk.CTkLabel(self.file_frame, text="No template selected. A default will be used.", anchor="w")
        self.template_label.grid(row=1, column=1, padx=10, pady=10, sticky="ew")

        # --- Middle Frame for Settings ---
        self.settings_frame = ctk.CTkFrame(self, corner_radius=10)
        self.settings_frame.grid(row=1, column=0, padx=20, pady=10, sticky="ew")
        self.settings_frame.grid_columnconfigure(1, weight=1)

        ctk.CTkLabel(self.settings_frame, text="Presentation Title:").grid(row=0, column=0, padx=(10,0), pady=10, sticky="w")
        self.ppt_title_entry = ctk.CTkEntry(self.settings_frame, placeholder_text="e.g., Website Assessment Report")
        self.ppt_title_entry.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        # --- Main Scrollable Frame for Column Mappings ---
        self.scrollable_frame = ctk.CTkScrollableFrame(self, label_text="Map Your Data Columns to PowerPoint Elements")
        self.scrollable_frame.grid(row=2, column=0, padx=20, pady=10, sticky="nsew")
        self.scrollable_frame.grid_columnconfigure(1, weight=1)

        # --- Bottom Frame for Actions ---
        self.bottom_frame = ctk.CTkFrame(self, corner_radius=10)
        self.bottom_frame.grid(row=3, column=0, padx=20, pady=(10, 20), sticky="ew")
        self.bottom_frame.grid_columnconfigure(0, weight=1)

        self.generate_button = ctk.CTkButton(self.bottom_frame, text="Generate Presentation", command=self.generate_presentation, state="disabled")
        self.generate_button.grid(row=0, column=0, padx=10, pady=10, sticky="ew")

        self.status_label = ctk.CTkLabel(self, text="Status: Ready. Please select a data file.", text_color="gray")
        self.status_label.grid(row=4, column=0, padx=20, pady=(0, 10), sticky="w")

    def load_data_file(self):
        """Loads data from Excel, CSV, or TSV files."""
        self.data_file_path = filedialog.askopenfilename(
            filetypes=(
                ("All Data Files", "*.xlsx *.xls *.csv *.tsv"),
                ("Excel files", "*.xlsx *.xls"),
                ("CSV files", "*.csv"),
                ("TSV files", "*.tsv"),
                ("All files", "*.*")
            )
        )
        if not self.data_file_path: return

        self.data_file_label.configure(text=os.path.basename(self.data_file_path))
        self.status_label.configure(text="Status: Reading data file...", text_color="white")
        self.update_idletasks()

        try:
            file_ext = os.path.splitext(self.data_file_path)[1].lower()
            header_row = self._find_header_row(self.data_file_path)

            if file_ext in ['.xlsx', '.xls']:
                self.dataframe = pd.read_excel(self.data_file_path, header=header_row)
            elif file_ext == '.csv':
                self.dataframe = pd.read_csv(self.data_file_path, header=header_row)
            elif file_ext == '.tsv':
                self.dataframe = pd.read_csv(self.data_file_path, sep='\t', header=header_row)
            else:
                raise ValueError("Unsupported file type.")

            self.dataframe.dropna(how='all', inplace=True)
            self.dataframe = self.dataframe.loc[:, ~self.dataframe.columns.str.contains('^Unnamed')]
            self.populate_column_mappings()
            self.generate_button.configure(state="normal")
            self.status_label.configure(text="Status: File loaded. Please map columns.", text_color="green")
        except Exception as e:
            messagebox.showerror("Error Loading File", f"Could not read the data file.\n\nError: {e}")
            self.status_label.configure(text="Status: Error loading file.", text_color="red")
            self.dataframe = None
            self.generate_button.configure(state="disabled")

    def _find_header_row(self, file_path):
        """Intelligently finds the header row in a data file."""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext in ['.xlsx', '.xls']:
                temp_df = pd.read_excel(file_path, header=None, nrows=10)
            else:
                temp_df = pd.read_csv(file_path, header=None, nrows=10, sep=None, engine='python')
            
            for i, row in temp_df.iterrows():
                if row.notna().sum() > len(row) / 2 and all(isinstance(x, str) for x in row if pd.notna(x)):
                    if row.nunique() >= len(row) / 2:
                        return i
        except Exception:
            return 0
        return 0

    def load_template_file(self):
        self.template_path = filedialog.askopenfilename(filetypes=(("PowerPoint templates", "*.pptx"),))
        if self.template_path:
            self.template_label.configure(text=os.path.basename(self.template_path))
        else:
            self.template_label.configure(text="No template selected. A default will be used.")

    def populate_column_mappings(self):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        self.column_widgets.clear()
        if self.dataframe is None: return

        for i, col_name in enumerate(self.dataframe.columns):
            label = ctk.CTkLabel(self.scrollable_frame, text=col_name, wraplength=250)
            label.grid(row=i, column=0, padx=10, pady=(5, 10), sticky="w")

            unique_count = self.dataframe[col_name].nunique()
            dtype = self.dataframe[col_name].dtype
            
            if "finding" in col_name.lower() or "summary" in col_name.lower():
                default_action = "Use for AI Summary"
            elif "risk" in col_name.lower() or "level" in col_name.lower():
                 default_action = "Create Pie Chart"
            elif "score" in col_name.lower() or dtype in ['int64', 'float64']:
                default_action = "Create Bar Chart"
            elif 0 < unique_count < 6:
                default_action = "Create Pie Chart"
            elif dtype == 'object' and unique_count > len(self.dataframe) * 0.8:
                 default_action = "Summarize as Bullet Points"
            else:
                default_action = "Ignore"

            dropdown_var = ctk.StringVar(value=default_action)
            dropdown = ctk.CTkOptionMenu(self.scrollable_frame, values=COLUMN_ACTIONS, variable=dropdown_var, width=250)
            dropdown.grid(row=i, column=1, padx=10, pady=(5, 10), sticky="e")
            self.column_widgets[col_name] = [label, dropdown]

    # --- PPTX HELPER FUNCTIONS ---
    def add_title_slide(self, prs, slide_title, slide_subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_subtitle

    def add_ai_summary_slide(self, prs, title, summary_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        tf.word_wrap = True

        disclaimer_text = ""
        if "(Note:" in summary_text:
            parts = summary_text.split("\n\n(Note:")
            summary_text = parts[0].strip()
            disclaimer_text = "(Note:" + parts[1].strip()

        lines = summary_text.strip().split('\n')
        for line in lines:
            line = line.strip()
            if not line: continue
            
            if line.startswith(('* ', '- ')):
                p = tf.add_paragraph()
                p.text = line[2:]
                p.level = 1; p.font.size = Pt(16)
            elif line.endswith(':'):
                p = tf.add_paragraph()
                p.text = line
                p.level = 0; p.font.bold = True; p.font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0; p.font.size = Pt(16)
        
        if disclaimer_text:
            tf.add_paragraph() 
            p = tf.add_paragraph()
            p.text = disclaimer_text
            p.level = 0; p.font.italic = True; p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(128, 128, 128)

    def add_section_header_slide(self, prs, title):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

    def add_bullet_point_slide(self, prs, title, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = str(point)
            p.level = 0

    def add_chart_slide(self, prs, title, chart_image_stream):
        """
        Adds a slide with a chart, ensuring the chart is centered within the
        content placeholder of the slide layout for proper alignment.
        """
        slide_layout = prs.slide_layouts[1] # Use 'Title and Content' layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        if len(slide.placeholders) < 2:
            slide.shapes.add_picture(chart_image_stream, Inches(2), Inches(2), width=Inches(6))
            return

        placeholder = slide.placeholders[1]
        
        chart_image_stream.seek(0)
        try:
            img = mpimg.imread(chart_image_stream, format='png')
            img_height_px, img_width_px, _ = img.shape
            image_aspect_ratio = float(img_width_px) / float(img_height_px)
        finally:
            chart_image_stream.seek(0)

        ph_width, ph_height = placeholder.width, placeholder.height
        placeholder_aspect_ratio = float(ph_width) / float(ph_height)

        if image_aspect_ratio > placeholder_aspect_ratio:
            new_width = ph_width
            new_height = new_width / image_aspect_ratio
        else:
            new_height = ph_height
            new_width = new_height * image_aspect_ratio

        left = placeholder.left + (ph_width - new_width) / 2
        top = placeholder.top + (ph_height - new_height) / 2
        
        slide.shapes.add_picture(chart_image_stream, left, top, width=new_width, height=new_height)

    def add_table_slide(self, prs, title, df_table):
        """
        Adds a slide with a table, placing it within the content
        placeholder of the slide for proper alignment.
        """
        slide_layout = prs.slide_layouts[1] # Use 'Title and Content' layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        if len(slide.placeholders) < 2:
            rows, cols = df_table.shape[0] + 1, df_table.shape[1]
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(5.5))
        else:
            placeholder = slide.placeholders[1]
            left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
            rows, cols = df_table.shape[0] + 1, df_table.shape[1]
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            placeholder.element.getparent().remove(placeholder.element)

        table = table_shape.table
        for c, col_name in enumerate(df_table.columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for r_idx, row_data in enumerate(df_table.itertuples(index=False), start=1):
            for c_idx, cell_data in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(cell_data)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                if (r_idx % 2) == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(220, 230, 241)

    async def get_ai_summary(self, data_for_summary):
        """Sends data to Gemini API and returns a textual summary."""
        self.status_label.configure(text="Status: Generating AI insights...", text_color="cyan")
        self.update_idletasks()
        
        MAX_ROWS_FOR_AI = 200
        disclaimer = ""
        if len(data_for_summary) > MAX_ROWS_FOR_AI:
            data_for_summary_sample = data_for_summary.sample(MAX_ROWS_FOR_AI, random_state=42)
            disclaimer = f"\n\n(Note: This summary is based on a random sample of {MAX_ROWS_FOR_AI} out of {len(data_for_summary)} total rows.)"
        else:
            data_for_summary_sample = data_for_summary

        data_string = data_for_summary_sample.to_csv(index=False)

        prompt = f"""
        As a professional data privacy and website security analyst, your task is to create an executive summary for a PowerPoint presentation.
        Based on the following audit data snippet:

        --- DATA ---
        {data_string}
        --- END DATA ---

        Please generate a summary that includes the following sections. Use bullet points for clarity.

        Executive Summary:
        - Provide a brief, high-level overview of the audit's main conclusions.

        Key Insights & Risk Areas:
        - Identify the most critical risks or recurring findings from the data.
        - Point out any significant trends (e.g., a specific type of vulnerability is common).

        Key Metrics:
        - Calculate and list important quantifiable metrics. For example: 'Total domains scanned: [count]', 'Domains with "High" risk: [count] ([percentage]%)', 'Average security score: [average]'.
        
        Your entire response should be formatted as plain text, ready to be copied into a presentation slide.
        """
        
        try:
            apiKey = os.environ.get("GOOGLE_API_KEY")
            if not apiKey:
                return "Error: GOOGLE_API_KEY environment variable not found. Please set it to use the AI summary feature."

            payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}]}
            apiUrl = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-pro:generateContent?key={apiKey}"
            
            response = await asyncio.to_thread(requests.post, apiUrl, json=payload, timeout=90)
            response.raise_for_status()
            
            result = response.json()
            
            if result.get('candidates'):
                summary_text = result['candidates'][0]['content']['parts'][0]['text']
                return summary_text + disclaimer
            else:
                error_info = result.get('promptFeedback', 'No content returned from API.')
                return f"Error: Could not generate AI summary. Details: {error_info}"

        except requests.exceptions.Timeout:
            return "An error occurred: The request to the AI service timed out."
        except requests.exceptions.RequestException as e:
            return f"An error occurred while contacting the AI service: {e}"
        except Exception as e:
            return f"An unexpected error occurred during AI summary generation: {e}"


    def generate_plots_for_df(self, prs, df_subset, group_title=""):
        mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}
        for col, action in mappings.items():
            if action not in ["Create Bar Chart", "Create Pie Chart"]: continue
            if df_subset[col].nunique() == 0 or df_subset[col].isnull().all(): continue

            plt.style.use('seaborn-v0_8-talk')
            fig, ax = plt.subplots(figsize=(10, 6))
            data_counts = df_subset[col].value_counts()
            chart_title = f"{group_title}: {col}" if group_title else f"Distribution of {col}"

            if action == "Create Bar Chart":
                data_counts.plot(kind='bar', ax=ax, color=plt.cm.viridis.colors)
                ax.set_ylabel("Count"); plt.xticks(rotation=45, ha='right')
            elif action == "Create Pie Chart":
                ax.pie(data_counts, labels=data_counts.index, autopct='%1.1f%%', startangle=140, colors=plt.cm.Pastel1.colors)
                ax.axis('equal')

            ax.set_title(chart_title, fontsize=16, pad=20)
            plt.tight_layout()

            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', dpi=200, bbox_inches='tight')
            plt.close(fig)
            img_stream.seek(0)
            self.add_chart_slide(prs, chart_title, img_stream)

        table_cols = [col for col, action in mappings.items() if action == "Include in Data Table"]
        if table_cols:
            table_df = df_subset[table_cols]
            table_title = f"{group_title}: Data Summary" if group_title else "Detailed Data Summary"
            self.add_table_slide(prs, table_title, table_df)

    async def generate_presentation_async(self):
        if self.dataframe is None:
            messagebox.showerror("Error", "No data has been loaded.")
            return

        self.generate_button.configure(state="disabled")
        self.status_label.configure(text="Status: Generating presentation...", text_color="white")
        self.update_idletasks()
        
        try:
            prs = Presentation(self.template_path) if self.template_path else Presentation()
            
            if self.template_path:
                # Remove all existing slides from the template file
                for i in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            
            if not self.template_path:
                prs.slide_width = Inches(16); prs.slide_height = Inches(9)

            ppt_title = self.ppt_title_entry.get() or "Data Analysis Report"
            subtitle_text = f"Source: {os.path.basename(self.data_file_path)}"
            self.add_title_slide(prs, ppt_title, subtitle_text)
            
            mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}
            
            ai_summary_cols = [col for col, action in mappings.items() if action == "Use for AI Summary"]
            if ai_summary_cols:
                summary_df = self.dataframe[ai_summary_cols]
                ai_summary_text = await self.get_ai_summary(summary_df)
                self.add_ai_summary_slide(prs, "AI-Powered Executive Summary", ai_summary_text)
                self.status_label.configure(text="Status: AI summary generated. Creating charts...", text_color="white")
                self.update_idletasks()

            summary_cols = [col for col, action in mappings.items() if action == "Summarize as Bullet Points"]
            if summary_cols:
                all_bullets = []
                for col in summary_cols:
                    bullets = self.dataframe[col].dropna().unique().tolist()
                    if bullets: all_bullets.extend(bullets)
                if all_bullets:
                    self.add_bullet_point_slide(prs, "Key Findings & Observations", all_bullets)

            grouping_col = next((col for col, action in mappings.items() if action == "Group Slides by this Column"), None)
            if grouping_col:
                unique_groups = self.dataframe[grouping_col].dropna().unique()
                for group in unique_groups:
                    self.add_section_header_slide(prs, f"Detailed Analysis for: {group}")
                    df_subset = self.dataframe[self.dataframe[grouping_col] == group]
                    self.generate_plots_for_df(prs, df_subset, group_title=str(group))
            else:
                self.generate_plots_for_df(prs, self.dataframe)

            self.add_section_header_slide(prs, "Thank You")

            save_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint Presentation", "*.pptx")],
                initialfile=f"{ppt_title.replace(' ', '_')}_Report.pptx",
                title="Save Presentation As"
            )

            if save_path:
                prs.save(save_path)
                self.status_label.configure(text=f"Status: Success! Saved to {os.path.basename(save_path)}", text_color="green")
                messagebox.showinfo("Success", f"Presentation generated successfully!\n\nSaved at: {save_path}")
            else:
                self.status_label.configure(text="Status: Save operation cancelled.", text_color="orange")

        except Exception as e:
            self.status_label.configure(text=f"Status: An error occurred during generation.", text_color="red")
            messagebox.showerror("Generation Error", f"An unexpected error occurred:\n\n{e}", detail=str(e))
            import traceback
            traceback.print_exc()
        finally:
            self.generate_button.configure(state="normal")


    def generate_presentation(self):
        try:
            asyncio.run(self.generate_presentation_async())
        except Exception as e:
             messagebox.showerror("Runtime Error", f"An error occurred while running the async operation:\n\n{e}")

if __name__ == "__main__":
    app = App()
    app.mainloop()
