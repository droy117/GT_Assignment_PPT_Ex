from dotenv import load_dotenv
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from tkinter import filedialog, messagebox
import customtkinter as ctk
import matplotlib.pyplot as plt
import os
import io
import asyncio
import requests
import json
import threading
from PIL import Image

load_dotenv()

# --- Constants and Theme Settings ---
APP_TITLE = "PPTEx"
WINDOW_SIZE = "1200x800"

COLUMN_ACTIONS = [
    "Ignore",
    "Create Bar Chart",
    "Create Pie Chart",
    "Create Histogram",
    "Create Line Chart",
]

API_KEY_FILE = "gemini_api_key.txt"

# --- Helper function to truncate long labels ---
def truncate_label(label, length=20):
    """Truncates a label if it's longer than the specified length."""
    label_str = str(label)
    if len(label_str) > length:
        return label_str[:length] + '...'
    return label_str

# --- Main Application Class ---
class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        # --- Window Setup ---
        self.title(APP_TITLE)
        self.geometry(WINDOW_SIZE)
        ctk.set_appearance_mode("Dark")
        ctk.set_default_color_theme("blue")

        # --- Class Variables ---
        self.dataframe = None
        self.original_df = None
        self.data_file_path = ""
        self.template_path = ""
        self.column_widgets = {}
        self.max_score_info = {}
        self.is_cancelling = False
        self.generation_thread = None

        # --- Settings Variables ---
        self.api_key = ctk.StringVar()
        self.ai_call_batch_size = ctk.IntVar(value=3)
        self.max_categories_for_charts = ctk.IntVar(value=9)

        # --- Set default template path if it exists ---
        default_template_file = './Templates/Template_1_default.pptx'
        if os.path.exists(default_template_file):
            self.template_path = default_template_file
        else:
            print(f"Default template not found at {default_template_file}. Will use a blank presentation.")

        # --- Main Layout Configuration ---
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(2, weight=1)

        # --- Top Frame for File Selection ---
        self.file_frame = ctk.CTkFrame(self, corner_radius=10)
        self.file_frame.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        self.file_frame.grid_columnconfigure(1, weight=1)

        self.browse_data_button = ctk.CTkButton(self.file_frame, text="Select Data File (Excel, CSV)", command=self.load_data_file)
        self.browse_data_button.grid(row=0, column=0, padx=10, pady=10)
        self.data_file_label = ctk.CTkLabel(self.file_frame, text="No data file selected", anchor="w")
        self.data_file_label.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        self.browse_template_button = ctk.CTkButton(self.file_frame, text="Select PowerPoint Template (Optional)", command=self.load_template_file)
        self.browse_template_button.grid(row=1, column=0, padx=10, pady=10)
        self.template_label = ctk.CTkLabel(self.file_frame, text=f"Template: {os.path.basename(self.template_path)}" if self.template_path else "No template selected. A default will be used.", anchor="w")
        self.template_label.grid(row=1, column=1, padx=10, pady=10, sticky="ew")

        # --- Middle Frame for Title and Settings Icon ---
        self.title_settings_frame = ctk.CTkFrame(self, corner_radius=10)
        self.title_settings_frame.grid(row=1, column=0, padx=20, pady=10, sticky="ew")
        self.title_settings_frame.grid_columnconfigure(1, weight=1)

        ctk.CTkLabel(self.title_settings_frame, text="Presentation Title:").grid(row=0, column=0, padx=(10,0), pady=10, sticky="w")
        self.ppt_title_entry = ctk.CTkEntry(self.title_settings_frame, placeholder_text="e.g., Website Assessment Report")
        self.ppt_title_entry.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        # --- Gear Icon for Settings ---
        self.settings_button = ctk.CTkButton(self.title_settings_frame, text="⚙️", width=32, command=self.toggle_settings_menu)
        self.settings_button.grid(row=0, column=2, padx=10, pady=10, sticky="e")
        self.title_settings_frame.grid_columnconfigure(2, weight=0)

        # --- Main Scrollable Frame for Column Mappings ---
        self.scrollable_frame = ctk.CTkScrollableFrame(self, label_text="Map Your Data Columns to PowerPoint Elements")
        self.scrollable_frame.grid(row=2, column=0, padx=20, pady=10, sticky="nsew")
        self.scrollable_frame.grid_columnconfigure(1, weight=1)

        # --- Bottom Frame for Actions ---
        self.bottom_frame = ctk.CTkFrame(self, corner_radius=10)
        self.bottom_frame.grid(row=3, column=0, padx=20, pady=(10, 0), sticky="ew")
        self.bottom_frame.grid_columnconfigure(0, weight=1)

        self.generate_button = ctk.CTkButton(self.bottom_frame, text="Generate Presentation", command=self.start_generation_thread, state="disabled")
        self.generate_button.grid(row=0, column=0, padx=10, pady=10, sticky="ew")

        # --- Status Log Textbox ---
        self.status_log = ctk.CTkTextbox(self, height=150, state="disabled", wrap="word")
        self.status_log.grid(row=4, column=0, padx=20, pady=(10, 20), sticky="ew")

        self.status_log.tag_config("INFO", foreground="gray")
        self.status_log.tag_config("SUCCESS", foreground="green")
        self.status_log.tag_config("ERROR", foreground="red")
        self.status_log.tag_config("WARN", foreground="orange")
        self.status_log.tag_config("CYAN", foreground="cyan")
        self.status_log.tag_config("WHITE", foreground="white")

        # --- Settings Menu (Initially Hidden) ---
        self.settings_menu_open = False
        self.settings_menu_frame = ctk.CTkFrame(self, corner_radius=15)
        self.create_settings_widgets()
        
        # --- Final Initialization ---
        self.load_api_key()
        self.check_api_key_and_update_ui()


    def create_settings_widgets(self):
        """Creates the widgets inside the settings menu frame."""
        self.settings_menu_frame.configure(border_width=3, border_color="#2196F3", corner_radius=0)

        settings_content_frame = ctk.CTkFrame(self.settings_menu_frame, fg_color="transparent")
        settings_content_frame.pack(padx=20, pady=20, fill="both", expand=True)
        
        # --- Close Button ---
        close_button = ctk.CTkButton(settings_content_frame, text="X", width=30, height=30, command=self.toggle_settings_menu)
        close_button.grid(row=0, column=2, sticky="ne", padx=(0,5), pady=(0,10))

        # --- API Key ---
        api_label = ctk.CTkLabel(settings_content_frame, text="Gemini API Key", font=ctk.CTkFont(size=16, weight="bold"))
        api_label.grid(row=1, column=0, columnspan=3, sticky="w", pady=(0, 5))
        
        self.api_key_entry = ctk.CTkEntry(settings_content_frame, textvariable=self.api_key, show="*", width=400)
        self.api_key_entry.grid(row=2, column=0, columnspan=3, sticky="ew")

        api_desc = ctk.CTkLabel(settings_content_frame, text="Your API key is stored locally in api_key.txt and is never shared. Get API key here: https://aistudio.google.com/app/apikey. For submission we have already added an API key in the gemini_api_key.txt file.",
                    font=ctk.CTkFont(slant="italic"), wraplength=450, justify="left")
        api_desc.grid(row=3, column=0, columnspan=3, sticky="w", pady=(0, 20))
        

        # --- Batch Size ---
        batch_label = ctk.CTkLabel(settings_content_frame, text="AI Call Batch Size", font=ctk.CTkFont(size=16, weight="bold"))
        batch_label.grid(row=4, column=0, columnspan=3, sticky="w", pady=(0, 5))

        self.batch_size_slider = ctk.CTkSlider(settings_content_frame, from_=1, to=10, number_of_steps=9, variable=self.ai_call_batch_size)
        self.batch_size_slider.grid(row=5, column=0, columnspan=2, sticky="ew", padx=(0, 10))
        self.batch_size_value_label = ctk.CTkLabel(settings_content_frame, textvariable=self.ai_call_batch_size)
        self.batch_size_value_label.grid(row=5, column=2, sticky="w")
        
        batch_desc = ctk.CTkLabel(settings_content_frame, text="Processes this many columns per single API call. A smaller number is slower but more reliable.",
                      font=ctk.CTkFont(slant="italic"), wraplength=450, justify="left")
        batch_desc.grid(row=6, column=0, columnspan=3, sticky="w", pady=(0, 20))

        # --- Max Categories ---
        cat_label = ctk.CTkLabel(settings_content_frame, text="Max Categories for Charts", font=ctk.CTkFont(size=16, weight="bold"))
        cat_label.grid(row=7, column=0, columnspan=3, sticky="w", pady=(0, 5))
        
        self.max_cat_slider = ctk.CTkSlider(settings_content_frame, from_=3, to=20, number_of_steps=17, variable=self.max_categories_for_charts)
        self.max_cat_slider.grid(row=8, column=0, columnspan=2, sticky="ew", padx=(0, 10))
        self.max_cat_value_label = ctk.CTkLabel(settings_content_frame, textvariable=self.max_categories_for_charts)
        self.max_cat_value_label.grid(row=8, column=2, sticky="w")
        
        cat_desc = ctk.CTkLabel(settings_content_frame, text="For pie/bar charts, categories with smaller counts will be grouped into an 'Others' slice. Larger counts will show more categories but might make the chart less readable.",
                    font=ctk.CTkFont(slant="italic"), wraplength=450, justify="left")
        cat_desc.grid(row=9, column=0, columnspan=3, sticky="w", pady=(0, 20))
        
        # --- Save Button ---
        save_button = ctk.CTkButton(settings_content_frame, text="Save and Close", command=self.save_and_close_settings)
        save_button.grid(row=10, column=0, columnspan=3, pady=(10,0))
    
    def animate_settings_menu(self, start_y, end_y):
        """Helper function to animate the sliding of the settings menu."""
        step = (end_y - start_y) / 10
        current_y = start_y
        
        def _animate():
            nonlocal current_y
            current_y += step
            if (step > 0 and current_y >= end_y) or (step < 0 and current_y <= end_y):
                self.settings_menu_frame.place_configure(rely=end_y)
                if end_y < 0: # If moved off-screen, forget it
                    self.settings_menu_frame.place_forget()
            else:
                self.settings_menu_frame.place_configure(rely=current_y)
                self.after(15, _animate)

        _animate()

    def toggle_settings_menu(self):
        """Opens or closes the settings menu with a sliding animation."""
        if self.settings_menu_open:
            self.animate_settings_menu(start_y=0.5, end_y=-1.0)
        else:
            self.settings_menu_frame.place(relx=0.5, rely=-1.0, anchor="center")
            self.animate_settings_menu(start_y=-1.0, end_y=0.5)
        
        self.settings_menu_open = not self.settings_menu_open


    def save_and_close_settings(self):
        """Saves the API key and closes the menu."""
        self.save_api_key()
        self.check_api_key_and_update_ui()
        self.toggle_settings_menu()

    def load_api_key(self):
        """Loads the API key from a local file."""
        if os.getenv("GEMINI_API_KEY"):
            self.api_key.set(os.getenv("GEMINI_API_KEY"))
            self.log_status("Loaded API key from environment variable.", "INFO")
            return
        try:
            if os.path.exists(API_KEY_FILE):
                with open(API_KEY_FILE, 'r') as f:
                    self.api_key.set(f.read().strip())
                self.log_status("Loaded API key from file.", "INFO")
        except Exception as e:
            self.log_status(f"Could not load API key from file: {e}", "ERROR")

    def save_api_key(self):
        """Saves the current API key to a local file."""
        try:
            with open(API_KEY_FILE, 'w') as f:
                f.write(self.api_key.get())
            self.log_status("API Key saved successfully.", "SUCCESS")
        except Exception as e:
            self.log_status(f"Failed to save API key: {e}", "ERROR")
            messagebox.showerror("Error", f"Could not save API key to {API_KEY_FILE}.\n\nError: {e}")

    def check_api_key_and_update_ui(self):
        """Checks if the API key exists and updates UI elements accordingly."""
        if not self.api_key.get():
            self.log_status("GEMINI_API_KEY is not set. Please add it in the settings menu (⚙️).", "ERROR")
            if self.dataframe is not None:
                self.generate_button.configure(state="disabled")
        else:
            self.log_status("API key is configured. Ready to generate.", "SUCCESS")
            if self.dataframe is not None:
                self.generate_button.configure(state="normal")
    
    def log_status(self, message, level="INFO"):
        self.status_log.configure(state="normal")
        self.status_log.insert("end", f"{message}\n", level.upper())
        self.status_log.configure(state="disabled")
        self.status_log.see("end")
        self.update_idletasks()

    def clean_data(self, df):
        self.log_status("Cleaning data...", "INFO")
        cleaned_df = df.copy()
        self.max_score_info = {}
        for col in cleaned_df.columns:
            if pd.api.types.is_numeric_dtype(cleaned_df[col].dtype):
                continue
            if cleaned_df[col].notna().any():
                col_as_str = cleaned_df[col].astype(str)
                pattern = r'^\s*([0-9.]+)\s*(out of|/)\s*([0-9.]+)'
                extracted_data = col_as_str.str.extract(pattern)
                if not extracted_data[0].isnull().all():
                    scores = pd.to_numeric(extracted_data[0], errors='coerce')
                    max_vals = pd.to_numeric(extracted_data[2], errors='coerce')
                    if scores.notna().any():
                        new_col_name = f"{col} (Score)"
                        cleaned_df[new_col_name] = scores
                        first_valid_max = max_vals.dropna().iloc[0] if max_vals.notna().any() else None
                        if first_valid_max:
                            self.max_score_info[new_col_name] = first_valid_max
                        cleaned_df.drop(columns=[col], inplace=True)
                        self.log_status(f"Cleaned and extracted scores from '{col}' into '{new_col_name}'.", "INFO")
                        continue
                if pd.api.types.is_object_dtype(cleaned_df[col].dtype):
                    cleaned_series = pd.to_numeric(
                        col_as_str.str.replace(',', '', regex=False),
                        errors='coerce'
                    )
                    if cleaned_series.notna().any():
                        cleaned_df[col] = cleaned_series
        self.log_status("Data cleaning complete.", "SUCCESS")
        return cleaned_df

    def load_data_file(self):
        self.data_file_path = filedialog.askopenfilename(
            filetypes=(("All Data Files", "*.xlsx *.xls *.csv *.tsv"),("Excel files", "*.xlsx *.xls"),("CSV files", "*.csv"),("TSV files", "*.tsv"),("All files", "*.*"))
        )
        if not self.data_file_path:
            self.log_status("File selection cancelled.", "WARN")
            return
        self.data_file_label.configure(text=os.path.basename(self.data_file_path))
        self.log_status(f"Reading data file: {os.path.basename(self.data_file_path)}", "WHITE")
        try:
            file_ext = os.path.splitext(self.data_file_path)[1].lower()
            header_row = self._find_header_row(self.data_file_path)
            self.log_status(f"Detected header at row {header_row + 1}.", "INFO")
            if file_ext in ['.xlsx', '.xls']:
                df = pd.read_excel(self.data_file_path, header=header_row)
            elif file_ext == '.csv':
                df = pd.read_csv(self.data_file_path, header=header_row)
            elif file_ext == '.tsv':
                df = pd.read_csv(self.data_file_path, sep='\t', header=header_row)
            else:
                raise ValueError("Unsupported file type.")
            df.dropna(how='all', inplace=True)
            df = df.loc[:, ~df.columns.str.contains('^Unnamed')]
            self.original_df = df.copy()
            self.dataframe = self.clean_data(df)
            self.populate_column_mappings()
            self.check_api_key_and_update_ui() # Check key after loading data
            self.log_status("File loaded successfully. Please map columns.", "SUCCESS")
        except Exception as e:
            messagebox.showerror("Error Loading File", f"Could not read the data file.\n\nError: {e}")
            self.log_status(f"Error loading file: {e}", "ERROR")
            self.dataframe = None
            self.generate_button.configure(state="disabled")

    def _find_header_row(self, file_path):
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext in ['.xlsx', '.xls']:
                temp_df = pd.read_excel(file_path, header=None, nrows=10)
            else:
                temp_df = pd.read_csv(file_path, header=None, nrows=10, sep=None, engine='python')
            for i, row in temp_df.iterrows():
                if row.notna().sum() > len(row) / 2 and all(isinstance(x, str) for x in row if pd.notna(x)):
                    if row.nunique() >= len(row) / 2:
                        return i
        except Exception:
            return 0
        return 0

    def load_template_file(self):
        self.template_path = filedialog.askopenfilename(filetypes=(("PowerPoint templates", "*.pptx"),))
        if self.template_path:
            self.template_label.configure(text=os.path.basename(self.template_path))
            self.log_status(f"Template selected: {os.path.basename(self.template_path)}", "INFO")
        else:
            self.template_label.configure(text="No template selected. A default will be used.")
            self.log_status("Template selection cancelled. Using default.", "WARN")

    def populate_column_mappings(self):
        """
        Populates the UI with column names and dropdowns for actions,
        with an improved intelligent default selection for the action.
        """
        # Clear existing widgets
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        self.column_widgets.clear()
        if self.dataframe is None: return

        self.log_status("Populating column mappings with improved intelligent defaults...", "INFO")

        for i, col_name in enumerate(self.dataframe.columns):
            label = ctk.CTkLabel(self.scrollable_frame, text=col_name, wraplength=250)
            label.grid(row=i, column=0, padx=10, pady=(5, 10), sticky="w")

            # --- New Recommendation Logic ---
            col = self.dataframe[col_name]
            unique_count = col.nunique()
            dtype = col.dtype
            default_action = "Ignore"  # Start with Ignore as the default

            # Rule 0: Ignore ID-like columns (new rule)
            id_keywords = ['id', 'sl.', 'sl ', 'serial', 'index', 'number', 'no.', 'no ', '#']
            if any(keyword in col_name.lower() for keyword in id_keywords):
                default_action = "Ignore"

            # Rule 1: Ignore date/time columns
            elif pd.api.types.is_datetime64_any_dtype(dtype):
                default_action = "Ignore"

            # Rule 2: Ignore columns with no or only one unique value
            elif unique_count <= 1:
                default_action = "Ignore"

            # Rule 3: Handle special 'Score' columns from the original logic
            elif col_name in self.max_score_info:
                default_action = "Create Histogram"

            # Rule 4: Handle boolean-like columns (True/False, Yes/No) -> Pie Chart
            elif unique_count == 2:
                default_action = "Create Pie Chart"

            # Rule 5: Handle numeric columns
            elif pd.api.types.is_numeric_dtype(dtype):
                # If there are a few unique numbers, treat as discrete categories -> Bar Chart
                if 2 < unique_count <= 10:
                    default_action = "Create Bar Chart"
                # For more unique values, it's likely continuous data -> Histogram
                else:  # unique_count > 10
                    default_action = "Create Histogram"

            # Rule 6: Handle categorical (object/string) columns
            elif pd.api.types.is_object_dtype(dtype):
                # If most values are unique (like IDs or comments), ignore them as summarization is removed
                if unique_count > len(col) * 0.8 and unique_count > 10:
                    default_action = "Ignore"
                # For a moderate number of categories -> Bar Chart
                elif unique_count > 7:  # User's suggestion
                    default_action = "Create Bar Chart"
                # For a small number of categories -> Pie Chart is good for proportions
                elif 2 < unique_count <= 7:
                    default_action = "Create Pie Chart"

            # Rule 7: Specific keyword check (from old logic) as a final override
            if "risk" in col_name.lower() or "level" in col_name.lower():
                default_action = "Create Pie Chart"
            # --- End of New Logic ---

            dropdown_var = ctk.StringVar(value=default_action)
            dropdown = ctk.CTkOptionMenu(self.scrollable_frame, values=COLUMN_ACTIONS, variable=dropdown_var, width=250)
            dropdown.grid(row=i, column=1, padx=10, pady=(5, 10), sticky="e")
            self.column_widgets[col_name] = [label, dropdown]
            
        self.log_status("Column mapping populated.", "INFO")

    def _find_relevant_columns_for_summary(self):
        """
        Identifies key columns for generating a conclusion by searching for keywords.
        """
        self.log_status("Identifying relevant columns for conclusion...", "INFO")
        if self.dataframe is None:
            return None

        # --- Keyword Definitions ---
        # Keywords to find the main subject of the report (e.g., the website list)
        website_keywords = ['website', 'domain', 'url', 'site', 'name']
        # Keywords to find metrics for the summary
        metric_keywords = ['compliance', 'score', 'risk', 'privacy', 'cookie', 'consent', 'gpc', 'geo', 'gap', 'level']

        website_col = None
        # Find the first column that matches website keywords
        for keyword in website_keywords:
            matching_cols = [col for col in self.dataframe.columns if keyword in col.lower()]
            if matching_cols:
                website_col = matching_cols[0]
                break # Stop after finding the first one

        metric_cols = []
        for keyword in metric_keywords:
            # Find columns that contain the keyword but are not the main website column
            matching_cols = [col for col in self.dataframe.columns if keyword in col.lower() and col != website_col]
            metric_cols.extend(matching_cols)

        # Get unique column names while preserving order
        metric_cols = list(dict.fromkeys(metric_cols))

        if not website_col and not metric_cols:
            self.log_status("Could not find any relevant columns for a conclusion slide.", "WARN")
            return None

        self.log_status(f"Found columns for summary. Website: '{website_col}', Metrics: {metric_cols}", "INFO")
        return {"website_col": website_col, "metric_cols": metric_cols}

    def _create_data_summary_for_conclusion(self, relevant_columns: dict):
        """
        Creates a text summary of key columns to be sent to the LLM.
        """
        if self.dataframe is None or not relevant_columns:
            return ""

        summary_parts = []
        website_col = relevant_columns.get("website_col")
        metric_cols = relevant_columns.get("metric_cols", [])

        total_rows = len(self.dataframe)
        summary_parts.append(f"The dataset contains {total_rows} total entries.")

        if website_col:
            summary_parts.append(f"\nThe analysis is for a list of websites, identified by the '{website_col}' column.")
            # Provide a small sample of website names for context
            if total_rows > 6:
                sample_websites = pd.concat([self.dataframe[website_col].head(3), self.dataframe[website_col].tail(3)]).unique()
                summary_parts.append(f"A sample of these websites includes: {', '.join(sample_websites)}.")
            else:
                summary_parts.append(f"The websites are: {', '.join(self.dataframe[website_col].unique())}.")

        summary_parts.append("\n\n--- Key Metric Summaries ---")

        for col in metric_cols:
            if col in self.dataframe.columns:
                series = self.dataframe[col].dropna()
                if pd.api.types.is_numeric_dtype(series.dtype):
                    summary_parts.append(f"\n\nColumn: '{col}' (Numeric Data)")
                    summary_parts.append(series.describe().to_string())
                else:
                    summary_parts.append(f"\n\nColumn: '{col}' (Categorical Data)")
                    summary_parts.append(series.value_counts().to_string())

        return "\n".join(summary_parts)

    async def get_ai_conclusion(self, summary_text: str):
        """
        Calls the Gemini API with a specific prompt to generate a conclusion.
        """
        self.log_status("Generating AI conclusion...", "CYAN")

        prompt = f"""
        You are an expert data analyst preparing the final slide for a PowerPoint presentation about a website compliance and assessment report.
        Based on the following data summary, write a concise, high-level conclusion for a non-technical audience.

        --- DATA SUMMARY ---
        {summary_text}
        --- END SUMMARY ---

        Instructions:
        1.  **Structure:** Write 6 to 7 insightful bullet points. Start each bullet point with '- '.
        2.  **Tone:** Professional and conclusive.
        3.  **Content:** Focus on the big picture. Do not just list the raw statistics from the summary. Instead, **interpret** them. For example, instead of saying "the average compliance score was 55.4," you could say, "The compliance scores indicate a significant performance gap, with many sites falling short of the desired standard." If you see many 'High Gap Quantity' entries, you might say, "A significant number of websites exhibit a high quantity of compliance gaps, pointing to systemic issues that need addressing."
        4.  **Context:** Refer to the subjects of the report as "the audited websites," "digital properties," or "the portfolio."
        5.  **Output:** Provide only the text for the bullet points. Do not include a title like "Conclusion" or any other surrounding text.

        Example Output:
        - The overall assessment reveals a significant variance in compliance maturity across the portfolio, with a notable cluster of websites having critical gaps that require immediate attention.
        - Automated compliance tools appear to be a key differentiator for achieving higher scores, as properties without such integrations consistently lag behind.
        - Geolocation-based rules and consent management are common areas of weakness, suggesting a need for a standardized implementation guide across all digital properties.
        """

        try:
            apiKey = self.api_key.get()
            if not apiKey:
                return "Error: GEMINI_API_KEY not found."

            payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}]}
            apiUrl = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={apiKey}"

            response = await asyncio.to_thread(requests.post, apiUrl, json=payload, timeout=90)
            response.raise_for_status()

            result = response.json()

            if result.get('candidates'):
                conclusion_text = result['candidates'][0]['content']['parts'][0]['text']
                self.log_status("Successfully generated AI conclusion.", "SUCCESS")
                return conclusion_text.strip()
            else:
                error_detail = f"Could not generate AI conclusion. Details: {result.get('promptFeedback', 'No content.')}"
                self.log_status(error_detail, "ERROR")
                return f"Error: {error_detail}"

        except Exception as e:
            error_msg = f"An unexpected error occurred during AI conclusion generation: {e}"
            self.log_status(error_msg, "ERROR")
            return error_msg

    def add_title_slide(self, prs, slide_title, slide_subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_subtitle

    def add_chart_and_insight_slide(self, prs, title, chart_image_stream, insight_text):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Handle chart placement
        placeholder_left = slide.placeholders[1]
        ph_element = placeholder_left.element
        ph_element.getparent().remove(ph_element)
        chart_image_stream.seek(0)
        slide.shapes.add_picture(chart_image_stream, placeholder_left.left, placeholder_left.top, width=placeholder_left.width)
        
        # Handle insight text
        placeholder_right = slide.placeholders[2]
        tf = placeholder_right.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        lines = insight_text.strip().replace('**', '').split('\n')
        first_paragraph = True
        
        for line in lines:
            line = line.strip()
            if not line: 
                continue
                
            if first_paragraph:
                # Use the existing first paragraph instead of adding a new one
                p = tf.paragraphs[0]
                first_paragraph = False
            else:
                # Add new paragraphs for subsequent lines
                p = tf.add_paragraph()
                
            if line.endswith(':'):
                p.text = line
                p.level = 0
                p.font.bold = True
                p.font.size = Pt(18)
            elif line.startswith(('* ', '- ')):
                p.text = line[2:]
                p.level = 1
                p.font.size = Pt(14)
            else:
                p.text = line
                p.level = 0
                p.font.size = Pt(14)
            
    def add_section_header_slide(self, prs, title):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title; title_shape.text = title
        placeholders_to_remove = [shp for shp in slide.shapes if shp.is_placeholder and shp.shape_id != title_shape.shape_id]
        for shape in placeholders_to_remove:
            sp_element = shape.element; sp_element.getparent().remove(sp_element)

    def add_bullet_point_slide(self, prs, title, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Get the content placeholder
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Add bullet points
        for i, point in enumerate(bullet_points):
            if i == 0:
                # Use the first paragraph that already exists
                p = tf.paragraphs[0]
            else:
                # Add new paragraphs for subsequent bullet points
                p = tf.add_paragraph()
            
            p.text = str(point)
            p.level = 0  # This creates the bullet point
            
            # Optional: Set font size for better appearance
            if hasattr(p, 'font'):
                p.font.size = Pt(16)

    def add_table_slide(self, prs, title, df_table):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        rows, cols = df_table.shape[0] + 1, df_table.shape[1]
        if len(slide.placeholders) < 2:
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(5.5))
        else:
            placeholder = slide.placeholders[1]
            table_shape = slide.shapes.add_table(rows, cols, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
            placeholder.element.getparent().remove(placeholder.element)
        table = table_shape.table
        for c, col_name in enumerate(df_table.columns):
            cell = table.cell(0, c); cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(79, 129, 189); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for r_idx, row_data in enumerate(df_table.itertuples(index=False), start=1):
            for c_idx, cell_data in enumerate(row_data):
                cell = table.cell(r_idx, c_idx); cell.text = str(cell_data); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                if (r_idx % 2) == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(220, 230, 241)

    def summarize_categorical_data(self, series: pd.Series, threshold: int):
        if not pd.api.types.is_object_dtype(series.dtype) and not pd.api.types.is_categorical_dtype(series.dtype):
            return series
        counts = series.value_counts()
        if len(counts) > threshold:
            top_n = counts.head(threshold).copy()
            other_count = counts.iloc[threshold:].sum()
            num_other_cats = len(counts) - threshold
            other_label = f"Others ({num_other_cats} categories)"
            top_n.loc[other_label] = other_count
            return top_n
        return counts

    async def get_ai_insights_for_batch(self, column_batch: list):
        batch_col_names = [col_name for col_name, _, _ in column_batch]
        self.log_status(f"Generating AI insights for: {', '.join(batch_col_names)}...", "CYAN")
        data_strings_concatenated = ""

        # --- MODIFICATION START ---
        # Added logic to explicitly pass the maximum score to the AI.
        for col_name, data, hint in column_batch:
            data_summary = ""
            header_info = "" # New variable to hold max score info

            # Check if the column is a score column and add its max value to the prompt context
            if col_name in self.max_score_info:
                max_score = self.max_score_info[col_name]
                header_info = f"Maximum Possible Score: {max_score}\n"

            if hint == 'counts':
                data_summary = data.to_string()
            elif hint == 'series':
                if pd.api.types.is_numeric_dtype(data.dtype):
                    data_summary = data.describe().to_string()
                else:
                    data_summary = data.value_counts().to_string()

            # Construct the final string for the column, including the max score info if available
            data_strings_concatenated += f"Column: '{col_name}'\n{header_info}---\n{data_summary}\n---\n\n"

        prompt = f"""
        As a data analyst, your task is to analyze data summaries for several columns from a dataset and generate insights for a PowerPoint presentation.
        The data represents charts that your user will see. For data with an 'Other' category, acknowledge that it represents a collection of smaller groups.

        --- DATA SUMMARIES ---
        {data_strings_concatenated}
        --- END DATA ---

        Based on this data, provide a concise analysis for EACH column. Your response MUST be a single, valid JSON object.
        The keys of the object must be the exact column names provided.
        The value for each key must be a string containing the analysis for that column, structured into two sections:

        **Summary Insight:**
        - Write a 3-4 sentence interpretation of what the data reveals, with each sentence as a separate bullet point starting with '-'. What is the main takeaway? Focus on the most significant findings, be descriptive but keep it short.

        **Key Metrics:**
        - List atleast 3-4 key, quantifiable metrics derived from the data summary.
        - **For categorical data (like for pie/bar charts, presented as category counts):**
            - You MUST calculate the total sum of all counts first.
            - Then, for the most significant categories (e.g., the top 3), list the category name, its raw count, AND its percentage of the total. Format it like: "Category Name: Count (Percentage%)".
        - **For numerical data (like for histograms, presented as descriptive statistics):**
            - Include key stats like Average, Max, Min, and total count.
        - **For score data (identified by a 'Maximum Possible Score' field in its summary):**
            - Use the provided 'Maximum Possible Score' as the correct denominator (the 'Y' in 'X out of Y').
            - Provide the Average, Max, and Min scores from the data summary.
            - When reporting these metrics, present them in an 'X out of Y' format. For the average score, also include the percentage of the maximum.

        Example JSON output format:
        {{
          "Risk Level": "**Summary Insight:**\\n- The analysis shows a significant portion of items are categorized as high risk.\\n\\n**Key Metrics:**\\n- High Risk: 75 (64.1%)\\n- Medium Risk: 25 (21.4%)",
          "Compliance Score (Score)": "**Summary Insight:**\\n- On average, the compliance score is high, but there's room for improvement to reach the maximum score.\\n\\n**Key Metrics:**\\n- Average Score: 95.5 out of 115 (83.0%)\\n- Maximum Score: 110 out of 115\\n- Minimum Score: 75 out of 115"
        }}
        **IMPORTANT:**Do not put any new line before 'Summary Insight'.
        Ensure the entire output is a single, valid JSON object and nothing else. Do not include markdown formatting like ```json in the final response.
        """
        try:    
            apiKey = self.api_key.get()
            if not apiKey: return {name: "Error: GEMINI_API_KEY not found. Please set it in the settings." for name in batch_col_names}
            payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}]}
            apiUrl = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={apiKey}"
            response = await asyncio.to_thread(requests.post, apiUrl, json=payload, timeout=120)
            response.raise_for_status()
            result = response.json()
            if result.get('candidates'):
                json_text = result['candidates'][0]['content']['parts'][0]['text']
                json_text = json_text.strip().lstrip('```json').rstrip('```')
                insights = json.loads(json_text)
                self.log_status(f"Successfully received AI insights for: {', '.join(insights.keys())}", "INFO")
                return insights
            else:
                error_detail = f"Could not generate AI summary. Details: {result.get('promptFeedback', 'No content.')}"
                self.log_status(error_detail, "ERROR")
                return {name: f"Error: {error_detail}" for name in batch_col_names}
        except json.JSONDecodeError as e:
            error_msg = f"AI response was not valid JSON: {e}"
            self.log_status(error_msg, "ERROR")
            return {name: error_msg for name in batch_col_names}
        except Exception as e:
            error_msg = f"An unexpected error occurred during AI insight generation: {e}"
            self.log_status(error_msg, "ERROR")
            return {name: error_msg for name in batch_col_names}

    async def generate_plots_for_df(self, prs, df_subset, group_title=""):
        import matplotlib
        matplotlib.use('Agg')
        mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}
        chart_actions = ["Create Bar Chart", "Create Pie Chart", "Create Histogram", "Create Line Chart"]
        
        # Precompute counts for categorical data
        precomputed_counts = {}
        for col, action in mappings.items():
            if self.is_cancelling: return
            if action in ["Create Bar Chart", "Create Pie Chart"]:
                if col in df_subset.columns and pd.api.types.is_object_dtype(df_subset[col].dtype):
                    if df_subset[col].nunique() > 0 and not df_subset[col].isnull().all():
                        summarized = self.summarize_categorical_data(df_subset[col], self.max_categories_for_charts.get())
                        precomputed_counts[col] = summarized
        
        # Prepare data for AI insights
        columns_for_ai = []
        for col, action in mappings.items():
            if self.is_cancelling: return
            if action in chart_actions and col in df_subset.columns and df_subset[col].nunique() > 0 and not df_subset[col].isnull().all():
                if col in precomputed_counts:
                    columns_for_ai.append((col, precomputed_counts[col], 'counts'))
                else:
                    columns_for_ai.append((col, df_subset[col], 'series'))
        
        # Get AI insights in batches
        all_insights = {}
        if columns_for_ai:
            if self.is_cancelling: return
            batches = [columns_for_ai[i:i + self.ai_call_batch_size.get()] for i in range(0, len(columns_for_ai), self.ai_call_batch_size.get())]
            tasks = [self.get_ai_insights_for_batch(batch) for batch in batches]
            results_list = await asyncio.gather(*tasks)
            for result_dict in results_list:
                if isinstance(result_dict, dict):
                    all_insights.update(result_dict)

        # Create chart generation tasks
        chart_tasks = []
        for col, action in mappings.items():
            if self.is_cancelling: return
            if action not in chart_actions: continue
            if col not in df_subset.columns or df_subset[col].nunique() == 0 or df_subset[col].isnull().all(): continue
            
            chart_tasks.append(self._create_chart_task(col, action, df_subset, precomputed_counts, all_insights, group_title))
        
        # Execute chart generation in parallel
        if chart_tasks:
            chart_results = await asyncio.gather(*chart_tasks, return_exceptions=True)
            
            # Add slides for successful charts
            for result in chart_results:
                if self.is_cancelling: return
                if isinstance(result, dict) and 'success' in result:
                    self.add_chart_and_insight_slide(prs, result['title'], result['img_stream'], result['insight'])
                    self.log_status(f"Successfully added slide for '{result['col']}'.", "INFO")
                elif isinstance(result, Exception):
                    self.log_status(f"Chart generation failed: {result}", "ERROR")

    async def _create_chart_task(self, col, action, df_subset, precomputed_counts, all_insights, group_title):
        """Creates a single chart and returns the data needed for slide creation."""
        import matplotlib.pyplot as plt
        import io
        
        self.log_status(f"Creating chart for '{col}'...", "WHITE")
        
        try:
            plt.style.use('seaborn-v0_8-talk')
        except OSError:
            plt.style.use('seaborn-talk')
        
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # --- Start of Addition ---
        # This block sets the font sizes for various chart elements.
        plt.rcParams.update({
            'font.size': 14,           # Set the global default font size
            'axes.titlesize': 16,      # Axis title
            'axes.labelsize': 14,      # Axis labels
            'xtick.labelsize': 12,     # X tick labels
            'ytick.labelsize': 12,     # Y tick labels
            'legend.fontsize': 12,     # Legend font
            'figure.titlesize': 18     # Figure title
        })
        # --- End of Addition ---

        chart_title = f"{group_title}: {col}" if group_title else f"Analysis of {col}"
        
        try:
            if action == "Create Bar Chart":
                if col in precomputed_counts:
                    data_counts = precomputed_counts[col]
                    data_counts.index = data_counts.index.map(truncate_label)
                    data_counts.plot(kind='bar', ax=ax, color=plt.cm.viridis.colors)
                    ax.set_ylabel("Count")
                elif col in self.max_score_info:
                    max_y_val = self.max_score_info[col]
                    df_subset[col].plot(kind='bar', ax=ax, color=plt.cm.viridis(df_subset[col].values / max_y_val))
                    ax.set_ylabel("Score")
                    ax.set_xlabel("Data Point Index")
                    ax.set_ylim(0, max_y_val * 1.05)
                    chart_title = f"{group_title}: {col.replace(' (Score)', '')}"
                else:
                    df_subset[col].value_counts().plot(kind='bar', ax=ax, color=plt.cm.viridis.colors)
                    ax.set_ylabel("Count")
                plt.xticks(rotation=45, ha='right')
                
            elif action == "Create Pie Chart":
                data_counts = precomputed_counts.get(col, df_subset[col].value_counts())
                truncated_labels = data_counts.index.map(truncate_label)
                # The 'textprops' argument is added here to control the font size inside the pie chart
                ax.pie(data_counts, labels=truncated_labels, autopct='%1.1f%%', startangle=140, 
                    colors=plt.cm.Pastel1.colors, textprops={'fontsize': 14})
                ax.axis('equal')
                
            elif action == "Create Histogram":
                if pd.api.types.is_numeric_dtype(df_subset[col]):
                    df_subset[col].dropna().plot(kind='hist', ax=ax, bins=15, color='skyblue', ec='black')
                    ax.set_ylabel("Frequency")
                    ax.set_xlabel(col)
                else:
                    plt.close(fig)
                    raise ValueError(f"Cannot create histogram for non-numeric column: '{col}'")
                    
            elif action == "Create Line Chart":
                if pd.api.types.is_numeric_dtype(df_subset[col]):
                    ax.plot(df_subset.index, df_subset[col], marker='o', linestyle='-')
                    ax.set_ylabel(col)
                    ax.set_xlabel("Index")
                else:
                    plt.close(fig)
                    raise ValueError(f"Cannot create line chart for non-numeric column: '{col}'")
            
            plt.tight_layout()
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', dpi=200, bbox_inches='tight')
            plt.close(fig)
            img_stream.seek(0)
            
            insight_text = all_insights.get(col, "AI insight could not be generated for this chart.")
            
            return {
                'success': True,
                'col': col,
                'title': chart_title,
                'img_stream': img_stream,
                'insight': insight_text
            }
            
        except Exception as e:
            plt.close(fig)
            raise Exception(f"Failed to create chart for '{col}': {e}")
        
    async def generate_presentation_async(self):
        if self.dataframe is None:
            messagebox.showerror("Error", "No data has been loaded.")
            return
        self.log_status("Starting presentation generation...", "WHITE")
        try:
            prs = Presentation(self.template_path) if self.template_path and os.path.exists(self.template_path) else Presentation()
            if self.is_cancelling: return
            if self.template_path and os.path.exists(self.template_path):
                self.log_status(f"Using template: {os.path.basename(self.template_path)}", "INFO")
                for i in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId; prs.part.drop_rel(rId); del prs.slides._sldIdLst[i]
            else:
                self.log_status("No template found or specified. Using a blank 16:9 presentation.", "INFO")
                prs.slide_width = Inches(16); prs.slide_height = Inches(9)
            if self.is_cancelling: return
            ppt_title = self.ppt_title_entry.get() or "Website Assessment Report"
            subtitle_text = f"Source: {os.path.basename(self.data_file_path)}"
            self.add_title_slide(prs, ppt_title, subtitle_text)
            self.log_status("Added title slide.", "INFO")
            if self.is_cancelling: return
            mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}

            if self.is_cancelling: return
            
            self.log_status("Generating slides for the entire dataset...", "WHITE")
            await self.generate_plots_for_df(prs, self.dataframe)

            self.log_status("Attempting to generate conclusion slide...", "WHITE")
            # Automatically find the most relevant columns for a summary
            relevant_columns_for_summary = self._find_relevant_columns_for_summary()
            
            if relevant_columns_for_summary:
                # Create a text summary from the identified columns
                data_summary_text = self._create_data_summary_for_conclusion(relevant_columns_for_summary)
                
                if data_summary_text:
                    # Call the AI to get the conclusion text
                    conclusion_text = await self.get_ai_conclusion(data_summary_text)
                    
                    if conclusion_text and not conclusion_text.lower().startswith("error:"):
                        # Process the AI response into a list of bullet points
                        bullet_points = [point.strip() for point in conclusion_text.split('\n') if point.strip().startswith('- ')]
                        cleaned_bullets = [bp.lstrip('- ').strip() for bp in bullet_points]
                        
                        # Add the conclusion slide to the presentation
                        self.add_bullet_point_slide(prs, "Conclusion", cleaned_bullets)
                        self.log_status("Successfully added conclusion slide.", "INFO")
                    else:
                        self.log_status(f"Could not add conclusion slide: {conclusion_text}", "WARN")

            if self.is_cancelling:
                self.log_status("Generation cancelled before saving.", "WARN"); return
                
                
            self.add_section_header_slide(prs, "Thank You")
            save_path = filedialog.asksaveasfilename(
                defaultextension=".pptx", filetypes=[("PowerPoint Presentation", "*.pptx")],
                initialfile=f"{ppt_title.replace(' ', '_')}_Report.pptx", title="Save Presentation As"
            )
            if save_path:
                prs.save(save_path)
                self.log_status(f"Success! Presentation saved to {os.path.basename(save_path)}", "SUCCESS")
                messagebox.showinfo("Success", f"Presentation generated successfully!\n\nSaved at: {save_path}")
            else:
                self.log_status("Save operation cancelled by user.", "WARN")
        except Exception as e:
            error_message = f"An unexpected error occurred during generation: {e}"
            self.log_status(error_message, "ERROR")
            messagebox.showerror("Generation Error", error_message, detail=str(e))
            import traceback; traceback.print_exc()

    def start_generation_thread(self):
        if self.generation_thread and self.generation_thread.is_alive():
            self.log_status("Generation is already in progress.", "WARN"); return
        self.is_cancelling = False
        self.generate_button.configure(text="Cancel Generation", command=self.cancel_generation, fg_color="#D32F2F", hover_color="#B71C1C")
        self.generation_thread = threading.Thread(target=self._run_async_generation, daemon=True)
        self.generation_thread.start()

    def cancel_generation(self):
        if self.generation_thread and self.generation_thread.is_alive():
            self.log_status("Cancellation requested. Will stop after the current step.", "WARN")
            self.is_cancelling = True
            self.generate_button.configure(text="Cancelling...", state="disabled")

    def _reset_generate_button_state(self):
        self.generate_button.configure(
            text="Generate Presentation", command=self.start_generation_thread,
            state="normal" if self.dataframe is not None and self.api_key.get() else "disabled",
            fg_color=ctk.ThemeManager.theme["CTkButton"]["fg_color"],
            hover_color=ctk.ThemeManager.theme["CTkButton"]["hover_color"]
        )
        self.is_cancelling = False
        self.generation_thread = None

    def _run_async_generation(self):
        try:
            asyncio.run(self.generate_presentation_async())
        except Exception as e:
            self.log_status(f"A runtime error occurred in the generation thread: {e}", "ERROR")
            self.after(0, lambda: messagebox.showerror("Runtime Error", f"An unexpected error occurred during generation:\n\n{e}"))
            import traceback; traceback.print_exc()
        finally:
            self.after(0, self._reset_generate_button_state)

if __name__ == "__main__":
    app = App()
    app.mainloop()