import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from tkinter import filedialog, messagebox
import customtkinter as ctk
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
import os
import winreg
import io
import asyncio
import requests
import re # Imported for data cleaning
from dotenv import load_dotenv
load_dotenv()


def get_windows_theme():
    """
    Returns 'Light' or 'Dark' depending on the current Windows app theme.
    """
    try:
        registry = winreg.ConnectRegistry(None, winreg.HKEY_CURRENT_USER)
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Themes\Personalize"
        key = winreg.OpenKey(registry, key_path)
        value, _ = winreg.QueryValueEx(key, "AppsUseLightTheme")
        return "Light" if value == 1 else "Dark"
    except Exception:
        return "Light"  # Fallback default

# --- NEW: Helper function to truncate long labels ---
def truncate_label(label, length=20):
    """Truncates a label if it's longer than the specified length."""
    label_str = str(label)  # Ensure the label is a string
    if len(label_str) > length:
        return label_str[:length] + '...'
    return label_str

# --- Constants and Theme Settings ---
APP_TITLE = "PPTEx"
WINDOW_SIZE = "1200x800"
COLUMN_ACTIONS = [
    "Ignore",
    "Group Slides by this Column",
    "Summarize as Bullet Points",
    "Create Bar Chart",
    "Create Pie Chart",
    "Create Histogram",
    "Create Line Chart",
    "Include in Data Table"
]

# --- Main Application Class ---
class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        # --- Window Setup ---
        self.title(APP_TITLE)
        self.geometry(WINDOW_SIZE)
        ctk.set_appearance_mode("Dark") # Force Dark mode as per user's prompt
        ctk.set_default_color_theme("blue")

        # --- Class Variables ---
        self.dataframe = None
        self.original_df = None
        self.data_file_path = ""
        # Set the default template path
        self.template_path = os.path.join(os.path.dirname(__file__), "Template.pptx") 
        self.column_widgets = {}
        self.max_score_info = {}

        # --- Main Layout Configuration ---
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(2, weight=1)

        # --- Top Frame for File Selection ---
        self.file_frame = ctk.CTkFrame(self, corner_radius=10)
        self.file_frame.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        self.file_frame.grid_columnconfigure(1, weight=1)

        self.browse_data_button = ctk.CTkButton(self.file_frame, text="Select Data File (Excel, CSV)", command=self.load_data_file)
        self.browse_data_button.grid(row=0, column=0, padx=10, pady=10)
        self.data_file_label = ctk.CTkLabel(self.file_frame, text="No data file selected", anchor="w")
        self.data_file_label.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        self.browse_template_button = ctk.CTkButton(self.file_frame, text="Select PowerPoint Template (Optional)", command=self.load_template_file)
        self.browse_template_button.grid(row=1, column=0, padx=10, pady=10)
        # Update the label to reflect the default template
        self.template_label = ctk.CTkLabel(self.file_frame, text=f"Default template: {os.path.basename(self.template_path)}", anchor="w")
        self.template_label.grid(row=1, column=1, padx=10, pady=10, sticky="ew")
        # Check if the default template exists and update the label accordingly
        if not os.path.exists(self.template_path):
            self.template_path = "" # Clear path if default doesn't exist
            self.template_label.configure(text="No default template found. A default will be used.")


        # --- Middle Frame for Settings ---
        self.settings_frame = ctk.CTkFrame(self, corner_radius=10)
        self.settings_frame.grid(row=1, column=0, padx=20, pady=10, sticky="ew")
        self.settings_frame.grid_columnconfigure(1, weight=1)

        ctk.CTkLabel(self.settings_frame, text="Presentation Title:").grid(row=0, column=0, padx=(10,0), pady=10, sticky="w")
        self.ppt_title_entry = ctk.CTkEntry(self.settings_frame, placeholder_text="e.g., Website Assessment Report")
        self.ppt_title_entry.grid(row=0, column=1, padx=10, pady=10, sticky="ew")

        # --- Main Scrollable Frame for Column Mappings ---
        self.scrollable_frame = ctk.CTkScrollableFrame(self, label_text="Map Your Data Columns to PowerPoint Elements")
        self.scrollable_frame.grid(row=2, column=0, padx=20, pady=10, sticky="nsew")
        self.scrollable_frame.grid_columnconfigure(1, weight=1)

        # --- Bottom Frame for Actions ---
        self.bottom_frame = ctk.CTkFrame(self, corner_radius=10)
        self.bottom_frame.grid(row=3, column=0, padx=20, pady=(10, 20), sticky="ew")
        self.bottom_frame.grid_columnconfigure(0, weight=1)

        self.generate_button = ctk.CTkButton(self.bottom_frame, text="Generate Presentation", command=self.generate_presentation, state="disabled")
        self.generate_button.grid(row=0, column=0, padx=10, pady=10, sticky="ew")

        self.status_label = ctk.CTkLabel(self, text="Status: Ready. Please select a data file.", text_color="gray")
        self.status_label.grid(row=4, column=0, padx=20, pady=(0, 10), sticky="w")

    def clean_data(self, df):
        cleaned_df = df.copy()
        self.max_score_info = {}

        for col in cleaned_df.columns:
            if cleaned_df[col].dtype == 'object' and cleaned_df[col].notna().any():
                pattern = r'^\s*([0-9.]+)\s*(out of|/)\s*([0-9.]+)'
                if cleaned_df[col].str.contains(pattern, regex=True, na=False).any():
                    extracted_data = cleaned_df[col].str.extract(pattern)
                    scores = pd.to_numeric(extracted_data[0], errors='coerce')
                    max_vals = pd.to_numeric(extracted_data[2], errors='coerce')
                    if scores.notna().any():
                        new_col_name = f"{col} (Score)"
                        cleaned_df[new_col_name] = scores
                        first_valid_max = max_vals.dropna().iloc[0] if max_vals.notna().any() else None
                        if first_valid_max:
                            self.max_score_info[new_col_name] = first_valid_max
                        cleaned_df.drop(columns=[col], inplace=True)
        return cleaned_df

    def load_data_file(self):
        self.data_file_path = filedialog.askopenfilename(
            filetypes=(("All Data Files", "*.xlsx *.xls *.csv *.tsv"),("Excel files", "*.xlsx *.xls"),("CSV files", "*.csv"),("TSV files", "*.tsv"),("All files", "*.*"))
        )
        if not self.data_file_path: return
        self.data_file_label.configure(text=os.path.basename(self.data_file_path))
        self.status_label.configure(text="Status: Reading data file...", text_color="white")
        self.update_idletasks()
        try:
            file_ext = os.path.splitext(self.data_file_path)[1].lower()
            header_row = self._find_header_row(self.data_file_path)
            if file_ext in ['.xlsx', '.xls']:
                df = pd.read_excel(self.data_file_path, header=header_row)
            elif file_ext == '.csv':
                df = pd.read_csv(self.data_file_path, header=header_row)
            elif file_ext == '.tsv':
                df = pd.read_csv(self.data_file_path, sep='\t', header=header_row)
            else:
                raise ValueError("Unsupported file type.")
            df.dropna(how='all', inplace=True)
            df = df.loc[:, ~df.columns.str.contains('^Unnamed')]
            self.original_df = df.copy()
            self.dataframe = self.clean_data(df)
            self.populate_column_mappings()
            self.generate_button.configure(state="normal")
            self.status_label.configure(text="Status: File loaded. Please map columns.", text_color="green")
        except Exception as e:
            messagebox.showerror("Error Loading File", f"Could not read the data file.\n\nError: {e}")
            self.status_label.configure(text="Status: Error loading file.", text_color="red")
            self.dataframe = None
            self.generate_button.configure(state="disabled")

    def _find_header_row(self, file_path):
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext in ['.xlsx', '.xls']:
                temp_df = pd.read_excel(file_path, header=None, nrows=10)
            else:
                temp_df = pd.read_csv(file_path, header=None, nrows=10, sep=None, engine='python')
            for i, row in temp_df.iterrows():
                if row.notna().sum() > len(row) / 2 and all(isinstance(x, str) for x in row if pd.notna(x)):
                    if row.nunique() >= len(row) / 2:
                        return i
        except Exception:
            return 0
        return 0

    def load_template_file(self):
        self.template_path = filedialog.askopenfilename(filetypes=(("PowerPoint templates", "*.pptx"),))
        if self.template_path:
            self.template_label.configure(text=os.path.basename(self.template_path))
        else:
            # If user cancels, revert to default template (if it exists) or indicate none
            default_template_path = os.path.join(os.path.dirname(__file__), "Template.pptx")
            if os.path.exists(default_template_path):
                self.template_path = default_template_path
                self.template_label.configure(text=f"Default template: {os.path.basename(self.template_path)}")
            else:
                self.template_path = ""
                self.template_label.configure(text="No template selected. A default will be used.")

    def populate_column_mappings(self):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        self.column_widgets.clear()
        if self.dataframe is None: return
        for i, col_name in enumerate(self.dataframe.columns):
            label = ctk.CTkLabel(self.scrollable_frame, text=col_name, wraplength=250)
            label.grid(row=i, column=0, padx=10, pady=(5, 10), sticky="w")
            unique_count = self.dataframe[col_name].nunique()
            dtype = self.dataframe[col_name].dtype
            col_lower = col_name.lower()
            if col_name in self.max_score_info:
                default_action = "Create Bar Chart"
            elif "risk" in col_lower or "level" in col_lower:
                    default_action = "Create Pie Chart"
            elif dtype in ['int64', 'float64']:
                default_action = "Create Histogram"
            elif 0 < unique_count < 6 and dtype == 'object':
                default_action = "Create Pie Chart"
            elif dtype == 'object' and unique_count > len(self.dataframe) * 0.8:
                    default_action = "Summarize as Bullet Points"
            else:
                default_action = "Ignore"
            dropdown_var = ctk.StringVar(value=default_action)
            dropdown = ctk.CTkOptionMenu(self.scrollable_frame, values=COLUMN_ACTIONS, variable=dropdown_var, width=250)
            dropdown.grid(row=i, column=1, padx=10, pady=(5, 10), sticky="e")
            self.column_widgets[col_name] = [label, dropdown]

    # --- PPTX HELPER FUNCTIONS ---
    def add_title_slide(self, prs, slide_title, slide_subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_subtitle

    def add_chart_and_insight_slide(self, prs, main_slide_title, chart_image_stream, insight_text, chart_subject):
        """
        Adds a slide with a chart on the left, AI insights on the right, and the chart's
        specific subject/title below the chart.
        """
        # Layout 3 is 'Two Content' in most standard templates.
        # This layout typically has a title placeholder at the top.
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set the main slide title (at the top of the slide)
        slide.shapes.title.text = main_slide_title

        # Get placeholders for chart and insights
        placeholder_left = slide.placeholders[1] # This is typically the left content placeholder
        placeholder_right = slide.placeholders[2] # This is typically the right content placeholder
        
        # --- Add Chart to Left Placeholder ---
        # Remove the placeholder text box before adding image to avoid overlapping text
        ph_element_left = placeholder_left.element
        ph_element_left.getparent().remove(ph_element_left)
        
        chart_image_stream.seek(0)
        # Add the picture using the placeholder's bounding box for position and size
        pic = slide.shapes.add_picture(chart_image_stream, placeholder_left.left, placeholder_left.top, 
                                        width=placeholder_left.width, height=placeholder_left.height)
        
        # --- Add Chart Subject below the Chart (manual text box) ---
        # Calculate position for the chart subject text box.
        # It will be below the chart placeholder area.
        chart_subject_top = placeholder_left.top + placeholder_left.height + Inches(0.2) # 0.2 inches padding
        chart_subject_left = placeholder_left.left
        chart_subject_width = placeholder_left.width
        chart_subject_height = Inches(0.5) # Sufficient height for one line of text

        txBox_chart_subject = slide.shapes.add_textbox(chart_subject_left, chart_subject_top, 
                                                        chart_subject_width, chart_subject_height)
        tf_chart_subject = txBox_chart_subject.text_frame
        p_chart_subject = tf_chart_subject.paragraphs[0]
        p_chart_subject.text = f"Chart: {chart_subject}" # Prefix with "Chart: " for clarity
        p_chart_subject.font.size = Pt(12)
        p_chart_subject.font.bold = True
        p_chart_subject.font.color.rgb = RGBColor(100, 100, 100) # A subtle gray color
        p_chart_subject.alignment = MSO_ANCHOR.TOP # Align text to the top of its box

        # --- Add AI Insight to Right Placeholder ---
        # Remove the placeholder text box for insights as well
        ph_element_right = placeholder_right.element
        ph_element_right.getparent().remove(ph_element_right)

        txBox_insight = slide.shapes.add_textbox(placeholder_right.left, placeholder_right.top, 
                                                  placeholder_right.width, placeholder_right.height)
        tf_insight = txBox_insight.text_frame
        tf_insight.clear()
        tf_insight.word_wrap = True
        tf_insight.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        lines = insight_text.strip().replace('**', '').split('\n')
        for line in lines:
            line = line.strip()
            if not line: continue
            
            if line.endswith(':'):
                p = tf_insight.add_paragraph()
                p.text = line
                p.level = 0; p.font.bold = True; p.font.size = Pt(18)
            elif line.startswith(('* ', '- ')):
                p = tf_insight.add_paragraph()
                p.text = line[2:]
                p.level = 1; p.font.size = Pt(14)
            else:
                p = tf_insight.add_paragraph()
                p.text = line
                p.level = 0; p.font.size = Pt(14)


    def add_section_header_slide(self, prs, title):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        placeholders_to_remove = [shp for shp in slide.shapes if shp.is_placeholder and shp.shape_id != title_shape.shape_id]
        for shape in placeholders_to_remove:
            sp_element = shape.element
            sp_element.getparent().remove(sp_element)

    def add_bullet_point_slide(self, prs, title, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = str(point)
            p.level = 0

    def add_table_slide(self, prs, title, df_table):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        rows, cols = df_table.shape[0] + 1, df_table.shape[1]
        if len(slide.placeholders) < 2:
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(5.5))
        else:
            placeholder = slide.placeholders[1]
            table_shape = slide.shapes.add_table(rows, cols, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
            placeholder.element.getparent().remove(placeholder.element)
        table = table_shape.table
        for c, col_name in enumerate(df_table.columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for r_idx, row_data in enumerate(df_table.itertuples(index=False), start=1):
            for c_idx, cell_data in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(cell_data)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                if (r_idx % 2) == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(220, 230, 241)

    async def get_ai_insight_for_chart(self, column_data: pd.Series, column_name: str):
        """
        Generates targeted AI insights for a single data column.
        """
        
        data_string = ""
        # Create a summary of the data to send to the AI
        if pd.api.types.is_numeric_dtype(column_data):
            data_string = column_data.describe().to_string()
        else:
            data_string = column_data.value_counts().to_string()

        prompt = f"""
        As a data analyst, analyze the following data summary for a column named '{column_name}'.
        The data represents a chart that your user will see.

        --- DATA SUMMARY ---
        {data_string}
        --- END DATA ---

        Based on this data, provide a concise analysis for a PowerPoint slide. Structure your response into two sections using markdown-style headers:

        **Summary Insight:**
        - Write a 1-2 sentence interpretation of what this data reveals. What is the main takeaway? (e.g., "The compliance scores are mostly high, but there is a notable cluster of low-scoring domains requiring attention.")

        **Key Metrics:**
        - List 2-4 key, quantifiable metrics derived from the data.
        - For numerical data, include: Average, Max, Min, and number of entries (count).
        - For categorical data, include: The most frequent category and its count, and the least frequent category and its count.
        - Format metrics clearly, like: "- Average Score: 85.5"

        Your entire response must be plain text, ready for a slide.
        """
        try:
            apiKey = os.getenv("GEMINI_API_KEY")
            if not apiKey: return "Error: GEMINI_API_KEY environment variable not found."
            payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}]}
            apiUrl = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={apiKey}"
            
            # Use asyncio.to_thread for the synchronous requests.post call
            response = await asyncio.to_thread(requests.post, apiUrl, json=payload, timeout=90)
            response.raise_for_status()
            result = response.json()
            if result.get('candidates'):
                summary_text = result['candidates'][0]['content']['parts'][0]['text']
                return summary_text
            else:
                return f"Error: Could not generate AI summary. Details: {result.get('promptFeedback', 'No content.')}"
        except Exception as e:
            return f"An unexpected error occurred during AI insight generation for '{column_name}': {e}"

    async def generate_plots_for_df(self, prs, df_subset, group_title=""):
        """
        Creates charts and calls the AI for insights for each plot.
        Also truncates long labels on bar and pie charts.
        """
        mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}
        chart_actions = ["Create Bar Chart", "Create Pie Chart", "Create Histogram", "Create Line Chart"]
        
        chart_data_for_slides = []
        ai_tasks = []

        for col, action in mappings.items():
            if action not in chart_actions: continue
            if col not in df_subset.columns or df_subset[col].nunique() == 0 or df_subset[col].isnull().all(): continue

            plt.style.use('seaborn-v0_8-talk')
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Remove Matplotlib's default title, as we'll add it in PPTX
            ax.set_title('') 

            # --- Generate Chart Image ---
            if action == "Create Bar Chart":
                if col in self.max_score_info:
                    max_y_val = self.max_score_info[col]
                    df_subset[col].plot(kind='bar', ax=ax, color=plt.cm.viridis(df_subset[col].values / max_y_val))
                    ax.set_ylabel("Score"); ax.set_xlabel("Data Point Index")
                    ax.set_ylim(0, max_y_val * 1.05)
                else:
                    data_counts = df_subset[col].value_counts()
                    data_counts.index = data_counts.index.map(truncate_label)
                    data_counts.plot(kind='bar', ax=ax, color=plt.cm.viridis.colors)
                    ax.set_ylabel("Count")
                plt.xticks(rotation=45, ha='right')
            elif action == "Create Pie Chart":
                data_counts = df_subset[col].value_counts()
                truncated_labels = data_counts.index.map(truncate_label)
                ax.pie(data_counts, labels=truncated_labels, autopct='%1.1f%%', startangle=140, colors=plt.cm.Pastel1.colors)
                ax.axis('equal')
            elif action == "Create Histogram":
                if pd.api.types.is_numeric_dtype(df_subset[col]):
                    df_subset[col].plot(kind='hist', ax=ax, bins=15, color='skyblue', ec='black')
                    ax.set_ylabel("Frequency"); ax.set_xlabel(col)
                else: plt.close(fig); continue
            elif action == "Create Line Chart":
                if pd.api.types.is_numeric_dtype(df_subset[col]):
                    ax.plot(df_subset.index, df_subset[col], marker='o', linestyle='-')
                    ax.set_ylabel(col); ax.set_xlabel("Index")
                else: plt.close(fig); continue
            
            plt.tight_layout()
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', dpi=200, bbox_inches='tight')
            plt.close(fig)
            img_stream.seek(0)
            
            # Use 'col' as the chart_subject to display below the chart
            chart_data_for_slides.append({'main_title': f"Analysis of {col}", 'image_stream': img_stream, 'chart_subject': col})
            ai_tasks.append(self.get_ai_insight_for_chart(df_subset[col], col))
        
        if ai_tasks:
            self.status_label.configure(text="Status: Generating AI insights (this might take a moment)...", text_color="cyan")
            self.update_idletasks()
            
            insights = await asyncio.gather(*ai_tasks, return_exceptions=True)
            
            for i, chart_info in enumerate(chart_data_for_slides):
                insight_text = insights[i]
                if isinstance(insight_text, Exception):
                    insight_text = f"Error generating AI insight: {insight_text}"
                    messagebox.showwarning("AI Insight Error", f"Could not generate AI insight for {chart_info['main_title']}: {insights[i]}")

                # Pass all required arguments to add_chart_and_insight_slide
                self.add_chart_and_insight_slide(prs, chart_info['main_title'], chart_info['image_stream'], insight_text, chart_info['chart_subject'])

        table_cols = [col for col, action in mappings.items() if action == "Include in Data Table"]
        if table_cols:
            table_df = df_subset[table_cols]
            table_title = f"{group_title}: Data Summary" if group_title else "Detailed Data Summary"
            self.add_table_slide(prs, table_title, table_df)

    async def generate_presentation_async(self):
        if self.dataframe is None:
            messagebox.showerror("Error", "No data has been loaded.")
            return

        self.generate_button.configure(state="disabled")
        self.status_label.configure(text="Status: Generating presentation...", text_color="white")
        self.update_idletasks()
        
        try:
            # Check if the template_path is valid and exists, otherwise use default Presentation()
            if self.template_path and os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
                # Clear existing slides from the template if it's used
                for i in range(len(prs.slides) - 1, -1, -1):
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            else:
                prs = Presentation()
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)

            ppt_title = self.ppt_title_entry.get() or "Data Analysis Report"
            subtitle_text = f"Source: {os.path.basename(self.data_file_path)}"
            self.add_title_slide(prs, ppt_title, subtitle_text)
            
            mappings = {col: widgets[1].get() for col, widgets in self.column_widgets.items()}
            
            summary_cols = [col for col, action in mappings.items() if action == "Summarize as Bullet Points"]
            if summary_cols:
                all_bullets = []
                for col in summary_cols:
                    bullets = self.dataframe[col].dropna().unique().tolist()
                    if bullets: all_bullets.extend(bullets)
                if all_bullets:
                    self.add_bullet_point_slide(prs, "Key Findings & Observations", all_bullets)

            grouping_col = next((col for col, action in mappings.items() if action == "Group Slides by this Column"), None)
            if grouping_col:
                unique_groups = self.dataframe[grouping_col].dropna().unique()
                for group in unique_groups:
                    self.add_section_header_slide(prs, f"Detailed Analysis for: {group}")
                    df_subset = self.dataframe[self.dataframe[grouping_col] == group]
                    await self.generate_plots_for_df(prs, df_subset, group_title=str(group))
            else:
                await self.generate_plots_for_df(prs, self.dataframe)

            self.add_section_header_slide(prs, "Thank You")

            save_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint Presentation", "*.pptx")],
                initialfile=f"{ppt_title.replace(' ', '_')}_Report.pptx",
                title="Save Presentation As"
            )

            if save_path:
                prs.save(save_path)
                self.status_label.configure(text=f"Status: Success! Saved to {os.path.basename(save_path)}", text_color="green")
                messagebox.showinfo("Success", f"Presentation generated successfully!\n\nSaved at: {save_path}")
            else:
                self.status_label.configure(text="Status: Save operation cancelled.", text_color="orange")

        except Exception as e:
            self.status_label.configure(text=f"Status: An error occurred during generation.", text_color="red")
            messagebox.showerror("Generation Error", f"An unexpected error occurred:\n\n{e}", detail=str(e))
            import traceback
            traceback.print_exc()
        finally:
            self.generate_button.configure(state="normal")

    def generate_presentation(self):
        try:
            asyncio.run(self.generate_presentation_async())
        except RuntimeError as e:
            try:
                loop = asyncio.get_running_loop()
                if loop.is_running():
                    loop.create_task(self.generate_presentation_async())
                else:
                    asyncio.run(self.generate_presentation_async())
            except RuntimeError:
                asyncio.run(self.generate_presentation_async())
        except Exception as e:
            messagebox.showerror("Runtime Error", f"An error occurred while running the async operation:\n\n{e}")

if __name__ == "__main__":
    app = App()
    app.mainloop()
